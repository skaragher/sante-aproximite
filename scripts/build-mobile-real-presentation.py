from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = ROOT / "docs"
SHOTS_DIR = DOCS_DIR / "mobile-real-shots"
OUTPUT = DOCS_DIR / "Presentation_Mobile_Sante_Aproximite_Captures_Reelles.pptx"


COLORS = {
  "navy": RGBColor(20, 32, 70),
  "blue": RGBColor(38, 88, 214),
  "teal": RGBColor(14, 148, 136),
  "purple": RGBColor(111, 66, 193),
  "red": RGBColor(220, 38, 38),
  "text": RGBColor(31, 41, 55),
  "muted": RGBColor(100, 116, 139),
  "line": RGBColor(218, 226, 237),
  "bg": RGBColor(245, 247, 251),
  "white": RGBColor(255, 255, 255),
}


def add_full_bg(slide, prs, color):
  shape = slide.shapes.add_shape(
    MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
  )
  shape.fill.solid()
  shape.fill.fore_color.rgb = color
  shape.line.fill.background()
  return shape


def add_text(slide, left, top, width, height, text, size=20, bold=False, color=None, align=PP_ALIGN.LEFT):
  box = slide.shapes.add_textbox(left, top, width, height)
  frame = box.text_frame
  frame.word_wrap = True
  p = frame.paragraphs[0]
  p.alignment = align
  run = p.add_run()
  run.text = text
  run.font.size = Pt(size)
  run.font.bold = bold
  run.font.color.rgb = color or COLORS["text"]
  return box


def add_card(slide, left, top, width, height, fill=COLORS["white"], line=COLORS["line"], radius=None):
  shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
  shape.fill.solid()
  shape.fill.fore_color.rgb = fill
  shape.line.color.rgb = line
  shape.line.width = Pt(1.2)
  if radius is not None:
    shape.adjustments[0] = radius
  return shape


def add_title_band(slide, title, subtitle, accent):
  add_card(slide, Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.95), fill=COLORS["white"], line=COLORS["white"])
  add_card(slide, Inches(0.7), Inches(0.55), Inches(2.0), Inches(0.28), fill=accent, line=accent)
  add_text(slide, Inches(0.72), Inches(0.50), Inches(8.2), Inches(0.45), title, size=25, bold=True, color=COLORS["navy"])
  add_text(slide, Inches(0.72), Inches(0.86), Inches(9.2), Inches(0.3), subtitle, size=11.5, color=COLORS["muted"])


def add_image_slide(prs, title, subtitle, image_name, bullets, accent):
  slide = prs.slides.add_slide(prs.slide_layouts[6])
  add_full_bg(slide, prs, COLORS["bg"])
  add_title_band(slide, title, subtitle, accent)

  add_card(slide, Inches(0.55), Inches(1.55), Inches(6.2), Inches(5.35))
  image_path = SHOTS_DIR / image_name
  slide.shapes.add_picture(str(image_path), Inches(0.72), Inches(1.72), width=Inches(5.86), height=Inches(5.0))

  add_card(slide, Inches(7.0), Inches(1.55), Inches(5.35), Inches(5.35))
  add_text(slide, Inches(7.3), Inches(1.9), Inches(4.7), Inches(0.5), "Ce que montre cet ecran", size=18, bold=True, color=COLORS["navy"])

  y = 2.45
  for bullet in bullets:
    add_card(slide, Inches(7.32), Inches(y), Inches(0.22), Inches(0.22), fill=accent, line=accent)
    add_text(slide, Inches(7.6), Inches(y - 0.03), Inches(4.35), Inches(0.6), bullet, size=14, color=COLORS["text"])
    y += 0.78

  add_text(
    slide,
    Inches(7.3),
    Inches(5.85),
    Inches(4.6),
    Inches(0.7),
    "Capture reelle de l'application mobile en fonctionnement.",
    size=11.5,
    color=COLORS["muted"],
  )
  return slide


def build():
  prs = Presentation()
  prs.slide_width = Inches(13.333)
  prs.slide_height = Inches(7.5)

  slide = prs.slides.add_slide(prs.slide_layouts[6])
  add_full_bg(slide, prs, COLORS["navy"])
  add_card(slide, Inches(0.6), Inches(0.65), Inches(12.1), Inches(6.2), fill=COLORS["white"], line=COLORS["white"])
  add_text(slide, Inches(1.0), Inches(1.0), Inches(5.8), Inches(0.5), "Sante Aproximite", size=28, bold=True, color=COLORS["navy"])
  add_text(slide, Inches(1.0), Inches(1.45), Inches(5.8), Inches(0.8), "Presentation mobile attractive\navec captures d'ecran reelles", size=22, bold=True, color=COLORS["blue"])
  add_text(
    slide,
    Inches(1.0),
    Inches(2.35),
    Inches(4.8),
    Inches(2.2),
    "Une application mobile concue pour rapprocher les citoyens des services de sante, accelerer les signalements d'urgence et renforcer la confiance dans le systeme sanitaire.",
    size=16,
    color=COLORS["text"],
  )
  add_card(slide, Inches(1.0), Inches(5.2), Inches(2.3), Inches(0.45), fill=COLORS["teal"], line=COLORS["teal"])
  add_text(slide, Inches(1.18), Inches(5.22), Inches(2.0), Inches(0.25), "Solution mobile terrain", size=13, bold=True, color=COLORS["white"])
  slide.shapes.add_picture(str(SHOTS_DIR / "current-screen.png"), Inches(7.0), Inches(1.0), width=Inches(4.9), height=Inches(5.9))

  add_image_slide(
    prs,
    "Trouver rapidement un centre",
    "Recherche geolocalisee, navigation et information utile au citoyen",
    "current-screen.png",
    [
      "Carte mobile avec centres visibles autour de l'utilisateur.",
      "Recherche par nom ou par service de sante.",
      "Notation, satisfaction et acces direct a l'itineraire.",
      "Consultation possible depuis le terrain, en conditions reelles.",
    ],
    COLORS["teal"],
  )

  add_image_slide(
    prs,
    "Urgence sanitaire",
    "Signalement cible vers le bon service medical ou de secours",
    "urgence-sanitaire.png",
    [
      "Choix du service de prise en charge : SAMU ou sapeurs-pompiers.",
      "Typologie d'urgence adaptee au service selectionne.",
      "Ajout du lieu, de la description, de photos et de la position GPS.",
      "Flux pense pour reduire le temps d'alerte et faciliter la coordination.",
    ],
    COLORS["red"],
  )

  add_image_slide(
    prs,
    "Urgence securitaire",
    "Alerte mobile pour la police, la gendarmerie ou la protection civile",
    "urgence-securitaire.png",
    [
      "Adressage direct a l'autorite competente selon la nature de l'incident.",
      "Prise de photo et capture de position GPS sur le terrain.",
      "Historique des urgences deja transmises par l'usager.",
      "Un module utile pour la securite communautaire et la gestion de crise.",
    ],
    COLORS["purple"],
  )

  add_image_slide(
    prs,
    "Plaintes et retour citoyen",
    "Un canal simple pour remonter les difficultes vecues dans les centres",
    "plaintes.png",
    [
      "Plainte liee a un centre ou plainte generale.",
      "Recherche rapide d'un etablissement depuis la base locale.",
      "Sujets predefinis pour guider l'usager et normaliser les remontes.",
      "Un levier concret de redevabilite et d'amelioration continue.",
    ],
    COLORS["blue"],
  )

  slide = prs.slides.add_slide(prs.slide_layouts[6])
  add_full_bg(slide, prs, COLORS["bg"])
  add_title_band(
    slide,
    "Pourquoi cette application mobile est strategique",
    "Une solution simple a deployer et facile a adopter",
    COLORS["teal"],
  )
  points = [
    ("Acces", "Le citoyen localise rapidement le bon point de prise en charge."),
    ("Reactivite", "Les urgences remontent avec contexte, service cible et geolocalisation."),
    ("Confiance", "Les plaintes et evaluations structurent l'ecoute usager."),
    ("Pilotage", "Le mobile alimente un dispositif plus large de supervision sanitaire."),
  ]
  y = 1.9
  for label, body in points:
    add_card(slide, Inches(0.9), Inches(y), Inches(11.5), Inches(0.95))
    add_text(slide, Inches(1.15), Inches(y + 0.15), Inches(2.0), Inches(0.3), label, size=17, bold=True, color=COLORS["navy"])
    add_text(slide, Inches(2.45), Inches(y + 0.14), Inches(9.3), Inches(0.42), body, size=14, color=COLORS["text"])
    y += 1.08
  add_text(slide, Inches(0.9), Inches(6.55), Inches(10.5), Inches(0.35), "Captures reelles generees depuis l'application mobile connectee.", size=11.5, color=COLORS["muted"])

  prs.save(OUTPUT)
  print(f"PPTX cree: {OUTPUT}")


if __name__ == "__main__":
  build()
