from pathlib import Path
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt
from docx.enum.text import WD_ALIGN_PARAGRAPH


ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = ROOT / "docs"
SHOTS_DIR = DOCS_DIR / "presentation-assets"
DOCX_PATH = DOCS_DIR / "Dossier_Mobile_Sante_Aproximite.docx"
PPTX_PATH = DOCS_DIR / "Presentation_Mobile_Sante_Aproximite.pptx"


def get_font(size: int, bold: bool = False):
    candidates = [
        "arialbd.ttf" if bold else "arial.ttf",
        "calibrib.ttf" if bold else "calibri.ttf",
        "DejaVuSans-Bold.ttf" if bold else "DejaVuSans.ttf",
    ]
    for name in candidates:
        try:
            return ImageFont.truetype(name, size=size)
        except OSError:
            continue
    return ImageFont.load_default()


FONT_H1 = get_font(34, bold=True)
FONT_H2 = get_font(24, bold=True)
FONT_H3 = get_font(18, bold=True)
FONT_BODY = get_font(16)
FONT_SMALL = get_font(13)


def draw_round_rect(draw, box, radius, fill, outline=None, width=1):
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def draw_text(draw, xy, text, font, fill):
    draw.text(xy, text, font=font, fill=fill)


def draw_multiline(draw, xy, text, font, fill, max_width, line_gap=6):
    words = text.split()
    lines = []
    current = ""
    for word in words:
        test = word if not current else f"{current} {word}"
        if draw.textlength(test, font=font) <= max_width:
            current = test
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)

    x, y = xy
    for line in lines:
        draw.text((x, y), line, font=font, fill=fill)
        y += font.size + line_gap
    return y


def new_canvas(title: str, subtitle: str = ""):
    img = Image.new("RGB", (1600, 900), "#f4f7fb")
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, 0, 1600, 86), fill="#0c5a7a")
    draw.text((40, 22), title, font=FONT_H1, fill="white")
    if subtitle:
        draw.text((980, 28), subtitle, font=FONT_SMALL, fill="#d9eef5")
    draw.rounded_rectangle((30, 110, 1570, 860), radius=28, fill="white", outline="#d8e4ef", width=2)
    return img, draw


def draw_phone_mockup(draw, x, y, w, h, screen_title, accent="#0c5a7a"):
    draw_round_rect(draw, (x, y, x + w, y + h), 36, "#111827", "#0f172a", 2)
    draw_round_rect(draw, (x + 16, y + 16, x + w - 16, y + h - 16), 28, "#f8fbfe")
    draw.rounded_rectangle((x + w // 2 - 70, y + 10, x + w // 2 + 70, y + 24), radius=8, fill="#1f2937")
    draw.rectangle((x + 24, y + 46, x + w - 24, y + 98), fill=accent)
    draw.text((x + 38, y + 62), screen_title, font=FONT_H3, fill="white")
    return (x + 24, y + 110, x + w - 24, y + h - 28)


def save_overview():
    img, draw = new_canvas("Dashboard de pilotage", "Vue d'ensemble du systeme")
    stats = [
        ("1 284", "Centres enregistres", "#e7f2ff", "#1d4ed8"),
        ("312", "Centres proches", "#edfdf4", "#15803d"),
        ("845", "Utilisateurs geres", "#fff7ed", "#c2410c"),
        ("29", "Alertes non traitees", "#fef2f2", "#b91c1c"),
    ]
    x = 60
    for value, label, bg, fg in stats:
        draw_round_rect(draw, (x, 150, x + 340, 280), 22, bg, "#d9e4ec", 2)
        draw.text((x + 24, 182), value, font=get_font(38, True), fill=fg)
        draw.text((x + 24, 235), label, font=FONT_BODY, fill="#334155")
        x += 365

    draw_round_rect(draw, (60, 330, 760, 800), 22, "#f8fbfe", "#dce8f0", 2)
    draw.text((84, 360), "Tendances strategiques", font=FONT_H2, fill="#0f172a")
    bars = [68, 84, 52, 91, 73]
    colors = ["#1d4ed8", "#16a34a", "#0ea5e9", "#f59e0b", "#7c3aed"]
    for i, h in enumerate(bars):
        bx = 110 + i * 120
        draw.rounded_rectangle((bx, 700 - h * 4, bx + 66, 700), radius=10, fill=colors[i])
        draw.text((bx + 8, 715), f"T{i+1}", font=FONT_SMALL, fill="#475569")
    draw_multiline(
        draw,
        (84, 740),
        "Le tableau de bord consolide les centres, les utilisateurs, les urgences "
        "et les indicateurs de supervision dans une seule interface.",
        FONT_BODY,
        "#334155",
        620,
    )

    draw_round_rect(draw, (810, 330, 1530, 800), 22, "#f8fbfe", "#dce8f0", 2)
    draw.text((840, 360), "Actions rapides", font=FONT_H2, fill="#0f172a")
    actions = [
        "Actualiser les centres",
        "Actualiser les centres proches",
        "Voir les alertes",
        "Suivre les plaintes",
    ]
    ay = 430
    for item in actions:
        draw_round_rect(draw, (840, ay, 1450, ay + 70), 18, "#ffffff", "#cfdce6", 2)
        draw.text((870, ay + 22), item, font=FONT_BODY, fill="#1e293b")
        ay += 90
    draw_multiline(
        draw,
        (840, 720),
        "Cette vue parle directement aux ministeres, directions regionales et bailleurs : "
        "elle rend le systeme lisible en un coup d'oeil.",
        FONT_BODY,
        "#334155",
        620,
    )
    path = SHOTS_DIR / "01-dashboard-overview.png"
    img.save(path)
    return path


def save_mobile_nearby():
    img, draw = new_canvas("Application mobile", "Trouver rapidement le centre adapte depuis un smartphone")
    sx1 = draw_phone_mockup(draw, 120, 150, 360, 640, "Centres a proximite", "#0c5a7a")
    left, top, right, bottom = sx1
    draw_round_rect(draw, (left + 14, top + 8, left + 86, top + 42), 10, "#ffffff", "#cfdce6", 2)
    draw.text((left + 26, top + 18), "5 km", font=FONT_SMALL, fill="#334155")
    draw_round_rect(draw, (left + 98, top + 8, right - 72, top + 42), 10, "#ffffff", "#cfdce6", 2)
    draw.text((left + 112, top + 18), "Rechercher un service...", font=FONT_SMALL, fill="#94a3b8")
    draw_round_rect(draw, (right - 58, top + 8, right - 14, top + 42), 10, "#0c5a7a", "#0c5a7a", 2)
    draw.text((right - 45, top + 18), "OK", font=FONT_SMALL, fill="white")
    draw.rectangle((left + 14, top + 60, right - 14, top + 240), fill="#dff2ea")
    for ox, oy in [(left + 60, top + 110), (left + 190, top + 140), (left + 110, top + 190)]:
        draw.ellipse((ox, oy, ox + 22, oy + 22), fill="#1d4ed8", outline="white", width=3)
    draw.ellipse((left + 150, top + 155, left + 174, top + 179), fill="#0ea5e9", outline="white", width=3)
    cy = top + 260
    cards = [
        ("Hopital Regional A", "2.4 km", "Urgences • Imagerie • Laboratoire"),
        ("Centre Medical B", "4.1 km", "Consultation • Vaccination"),
        ("Maternite C", "6.8 km", "Maternite • Pediatrie"),
    ]
    for name, km, details in cards:
        draw_round_rect(draw, (left + 14, cy, right - 14, cy + 94), 14, "#ffffff", "#d8e4ef", 2)
        draw.text((left + 28, cy + 16), name, font=FONT_H3, fill="#0f172a")
        draw.text((right - 90, cy + 18), km, font=FONT_SMALL, fill="#0c5a7a")
        draw_multiline(draw, (left + 28, cy + 48), details, FONT_SMALL, "#475569", 230, 4)
        cy += 108

    sx2 = draw_phone_mockup(draw, 560, 150, 360, 640, "Itineraire & avis", "#0b7285")
    left, top, right, bottom = sx2
    draw.rectangle((left + 14, top + 10, right - 14, top + 220), fill="#e2f1e8")
    draw.line((left + 80, top + 170, right - 90, top + 80), fill="#0b7285", width=5)
    draw.ellipse((left + 72, top + 162, left + 96, top + 186), fill="#0ea5e9", outline="white", width=3)
    draw.ellipse((right - 104, top + 68, right - 80, top + 92), fill="#dc2626", outline="white", width=3)
    draw.text((left + 18, top + 236), "Hopital Regional A", font=FONT_H3, fill="#0f172a")
    draw.text((left + 18, top + 268), "Plateau: Urgences, Imagerie", font=FONT_SMALL, fill="#475569")
    draw.text((left + 18, top + 295), "4.3/5 • 87% satisfaits", font=FONT_SMALL, fill="#15803d")
    star_x = left + 18
    for i in range(5):
        active = i < 4
        draw_round_rect(draw, (star_x, top + 325, star_x + 34, top + 359), 17, "#0c5a7a" if active else "#ffffff", "#cfdce6", 2)
        draw.text((star_x + 10, top + 333), "★", font=FONT_SMALL, fill="white" if active else "#94a3b8")
        star_x += 40
    draw_round_rect(draw, (left + 18, top + 390, right - 18, top + 430), 10, "#0b7285", "#0b7285", 2)
    draw.text((left + 115, top + 401), "Lancer l'itineraire", font=FONT_SMALL, fill="white")

    draw_multiline(
        draw,
        (1000, 260),
        "La valeur mobile est immediate : l'usager localise un centre, compare les services, "
        "evalue la qualite et ouvre directement la navigation depuis son telephone.",
        FONT_BODY,
        "#334155",
        480,
    )
    draw_multiline(
        draw,
        (1000, 430),
        "Pour les autorites et bailleurs, cet usage mobile montre un impact concret sur l'acces "
        "aux soins, la reduction du temps de recherche et la satisfaction des patients.",
        FONT_BODY,
        "#334155",
        480,
    )
    path = SHOTS_DIR / "mobile-01-nearby.png"
    img.save(path)
    return path


def save_mobile_emergency():
    img, draw = new_canvas("Application mobile", "Signaler une urgence en quelques gestes")
    sx1 = draw_phone_mockup(draw, 100, 150, 360, 640, "Signaler une urgence", "#b91c1c")
    left, top, right, bottom = sx1
    services = [("SAMU", "#dc2626"), ("Pompiers", "#f97316"), ("Police", "#1d4ed8"), ("Protection", "#7c3aed")]
    sy = top + 10
    for i, (name, color) in enumerate(services):
        x = left + 14 + (i % 2) * 150
        y = sy + (i // 2) * 78
        draw_round_rect(draw, (x, y, x + 132, y + 60), 12, color if i == 0 else "#ffffff", color, 2)
        draw.text((x + 18, y + 16), name, font=FONT_SMALL, fill="white" if i == 0 else color)
    fy = top + 180
    labels = ["Telephone", "Type d'urgence", "Point de prise en charge", "Description"]
    heights = [36, 36, 36, 86]
    for label, hh in zip(labels, heights):
        draw.text((left + 18, fy - 18), label, font=FONT_SMALL, fill="#64748b")
        draw_round_rect(draw, (left + 18, fy, right - 18, fy + hh), 10, "#ffffff", "#cfdce6", 2)
        fy += hh + 34
    draw_round_rect(draw, (left + 18, fy, right - 18, fy + 48), 10, "#ecfdf5", "#86efac", 2)
    draw.text((left + 36, fy + 15), "GPS correct ~ 12 m", font=FONT_SMALL, fill="#15803d")
    draw_round_rect(draw, (left + 18, fy + 70, right - 18, fy + 118), 10, "#b91c1c", "#b91c1c", 2)
    draw.text((left + 74, fy + 85), "Envoyer au SAMU", font=FONT_SMALL, fill="white")

    sx2 = draw_phone_mockup(draw, 540, 150, 360, 640, "Suivi de mes alertes", "#0b7285")
    left, top, right, bottom = sx2
    cards = [
        ("Accident de route", "Equipe en route"),
        ("Detresse respiratoire", "Equipe sur site"),
        ("Evacuation obstetricale", "Intervention terminee"),
    ]
    cy = top + 16
    for title, status in cards:
        draw_round_rect(draw, (left + 18, cy, right - 18, cy + 110), 14, "#ffffff", "#d8e4ef", 2)
        draw.text((left + 34, cy + 18), title, font=FONT_H3, fill="#0f172a")
        draw_round_rect(draw, (left + 34, cy + 50, right - 34, cy + 78), 10, "#eff6ff", "#bfdbfe", 1)
        draw.text((left + 48, cy + 58), status, font=FONT_SMALL, fill="#1d4ed8")
        draw.text((left + 34, cy + 86), "Service: SAMU", font=FONT_SMALL, fill="#475569")
        cy += 126

    draw_multiline(
        draw,
        (980, 260),
        "Le mobile n'est pas seulement un canal d'envoi. Il devient un outil de sauvetage : "
        "service cible, telephone, GPS, description, photos et retour de statut.",
        FONT_BODY,
        "#334155",
        500,
    )
    draw_multiline(
        draw,
        (980, 430),
        "Cette experience smartphone parle fortement aux institutions car elle montre un usage terrain "
        "immediat pour le citoyen et une meilleure coordination des reponses d'urgence.",
        FONT_BODY,
        "#334155",
        500,
    )
    path = SHOTS_DIR / "mobile-02-emergency.png"
    img.save(path)
    return path


def save_mobile_complaints():
    img, draw = new_canvas("Application mobile", "Plainte, suivi et retour usager depuis le telephone")
    sx1 = draw_phone_mockup(draw, 110, 150, 360, 640, "Poser une plainte", "#1d4ed8")
    left, top, right, bottom = sx1
    draw_round_rect(draw, (left + 18, top + 18, left + 150, top + 54), 12, "#1d4ed8", "#1d4ed8", 2)
    draw.text((left + 36, top + 30), "Avec un centre", font=FONT_SMALL, fill="white")
    draw_round_rect(draw, (left + 162, top + 18, right - 18, top + 54), 12, "#ffffff", "#cfdce6", 2)
    draw.text((left + 182, top + 30), "Generale", font=FONT_SMALL, fill="#475569")
    fy = top + 92
    for label, hh in [("Centre de sante", 82), ("Sujet", 70), ("Description", 130)]:
        draw.text((left + 18, fy - 18), label, font=FONT_SMALL, fill="#64748b")
        draw_round_rect(draw, (left + 18, fy, right - 18, fy + hh), 10, "#ffffff", "#cfdce6", 2)
        fy += hh + 30
    draw_round_rect(draw, (left + 18, fy, right - 18, fy + 48), 10, "#1d4ed8", "#1d4ed8", 2)
    draw.text((left + 92, fy + 15), "Envoyer la plainte", font=FONT_SMALL, fill="white")

    sx2 = draw_phone_mockup(draw, 550, 150, 360, 640, "Mes plaintes", "#0c5a7a")
    left, top, right, bottom = sx2
    draw_round_rect(draw, (left + 18, top + 18, left + 152, top + 52), 12, "#0c5a7a", "#0c5a7a", 2)
    draw.text((left + 42, top + 28), "En cours (3)", font=FONT_SMALL, fill="white")
    draw_round_rect(draw, (left + 165, top + 18, right - 18, top + 52), 12, "#ffffff", "#cfdce6", 2)
    draw.text((left + 192, top + 28), "Resolues (6)", font=FONT_SMALL, fill="#475569")
    cy = top + 80
    for title, status, color in [
        ("Temps d'attente", "EN COURS", "#f59e0b"),
        ("Accueil", "RESOLUE", "#16a34a"),
        ("Tarification", "REJETEE", "#dc2626"),
    ]:
        draw_round_rect(draw, (left + 18, cy, right - 18, cy + 126), 14, "#ffffff", "#d8e4ef", 2)
        draw.text((left + 34, cy + 18), title, font=FONT_H3, fill="#0f172a")
        draw_round_rect(draw, (right - 120, cy + 16, right - 34, cy + 42), 10, "#ffffff", color, 2)
        draw.text((right - 108, cy + 23), status, font=FONT_SMALL, fill=color)
        draw.text((left + 34, cy + 56), "Hopital General • 23/04/2026", font=FONT_SMALL, fill="#64748b")
        draw_multiline(draw, (left + 34, cy + 82), "Le suivi mobile permet de voir les reponses et de donner un retour final.", FONT_SMALL, "#475569", 250, 4)
        cy += 142

    draw_multiline(
        draw,
        (980, 300),
        "Le parcours mobile de plainte renforce la redevabilite publique : depose simple, "
        "historique visible, statuts lisibles et feedback final de l'usager.",
        FONT_BODY,
        "#334155",
        500,
    )
    path = SHOTS_DIR / "mobile-03-complaints.png"
    img.save(path)
    return path


def save_mobile_chef():
    img, draw = new_canvas("Application mobile", "Le responsable de centre gere aussi depuis son telephone")
    sx1 = draw_phone_mockup(draw, 120, 150, 360, 640, "Mon centre de sante", "#0b7285")
    left, top, right, bottom = sx1
    draw.text((left + 18, top + 20), "Statut", font=FONT_SMALL, fill="#64748b")
    draw_round_rect(draw, (left + 78, top + 14, left + 190, top + 42), 10, "#ecfdf5", "#86efac", 2)
    draw.text((left + 96, top + 21), "APPROUVE", font=FONT_SMALL, fill="#15803d")
    fy = top + 70
    for label in ["Nom du centre", "Adresse", "Region", "District", "Services", "Latitude / Longitude"]:
        draw.text((left + 18, fy - 18), label, font=FONT_SMALL, fill="#64748b")
        draw_round_rect(draw, (left + 18, fy, right - 18, fy + 36), 10, "#ffffff", "#cfdce6", 2)
        fy += 58
    draw_round_rect(draw, (left + 18, fy, right - 18, fy + 44), 10, "#0b7285", "#0b7285", 2)
    draw.text((left + 80, fy + 13), "Mettre a jour", font=FONT_SMALL, fill="white")

    sx2 = draw_phone_mockup(draw, 560, 150, 360, 640, "Suivi du centre", "#1d4ed8")
    left, top, right, bottom = sx2
    stats = [("4.3", "Note"), ("87%", "Satisfaction"), ("12", "Plaintes")]
    x = left + 18
    for val, label in stats:
        draw_round_rect(draw, (x, top + 18, x + 92, top + 88), 12, "#eff6ff", "#bfdbfe", 2)
        draw.text((x + 18, top + 34), val, font=FONT_H3, fill="#1d4ed8")
        draw.text((x + 18, top + 60), label, font=FONT_SMALL, fill="#64748b")
        x += 102
    cy = top + 118
    for title in ["Temps d'attente aux urgences", "Disponibilite des medicaments"]:
        draw_round_rect(draw, (left + 18, cy, right - 18, cy + 150), 14, "#ffffff", "#d8e4ef", 2)
        draw.text((left + 34, cy + 18), title, font=FONT_H3, fill="#0f172a")
        draw.text((left + 34, cy + 48), "Ajouter une explication...", font=FONT_SMALL, fill="#94a3b8")
        draw_round_rect(draw, (left + 34, cy + 74, right - 34, cy + 108), 10, "#ffffff", "#cfdce6", 2)
        draw_round_rect(draw, (left + 34, cy + 116, right - 34, cy + 144), 10, "#ffffff", "#cfdce6", 2)
        draw.text((left + 82, cy + 123), "Ajouter explication", font=FONT_SMALL, fill="#475569")
        cy += 166

    draw_multiline(
        draw,
        (980, 290),
        "Le mobile ne sert pas qu'au citoyen. Le chef d'etablissement peut renseigner son centre, "
        "mettre a jour les services, suivre les plaintes et repondre rapidement depuis le terrain.",
        FONT_BODY,
        "#334155",
        500,
    )
    path = SHOTS_DIR / "mobile-04-chef.png"
    img.save(path)
    return path


def save_nearby():
    img, draw = new_canvas("Orientation vers les soins", "Recherche des centres a proximite")
    draw_round_rect(draw, (60, 150, 1540, 235), 20, "#f8fbfe", "#d9e4ef", 2)
    toolbar = [
        ("Rayon (km)", "20"),
        ("Recherche", "nom, type, service"),
        ("Position", "GPS actif"),
    ]
    tx = 86
    for label, value in toolbar:
        draw.text((tx, 165), label, font=FONT_SMALL, fill="#64748b")
        draw_round_rect(draw, (tx, 188, tx + 250, 222), 10, "#ffffff", "#cfdce6", 2)
        draw.text((tx + 10, 194), value, font=FONT_SMALL, fill="#0f172a")
        tx += 290
    draw_round_rect(draw, (1110, 182, 1290, 222), 10, "#dbeafe", "#93c5fd", 2)
    draw.text((1144, 193), "Ma position", font=FONT_SMALL, fill="#1d4ed8")
    draw_round_rect(draw, (1310, 182, 1460, 222), 10, "#0c5a7a", "#0c5a7a", 2)
    draw.text((1352, 193), "Rechercher", font=FONT_SMALL, fill="white")

    draw_round_rect(draw, (60, 270, 860, 800), 22, "#eaf6fb", "#dce8f0", 2)
    draw.text((88, 295), "Carte sanitaire", font=FONT_H2, fill="#0f172a")
    draw.rectangle((95, 345, 825, 765), fill="#dff2ea")
    for x in range(120, 820, 110):
        draw.line((x, 360, x, 750), fill="#c8ded3", width=2)
    for y in range(380, 760, 90):
        draw.line((110, y, 810, y), fill="#c8ded3", width=2)
    for ox, oy, color in [(260, 470, "#1d4ed8"), (520, 540, "#16a34a"), (690, 430, "#f59e0b"), (390, 650, "#1d4ed8")]:
        draw.ellipse((ox, oy, ox + 26, oy + 26), fill=color, outline="white", width=3)
    draw.ellipse((430, 500, 460, 530), fill="#0ea5e9", outline="white", width=3)
    draw.line((445, 515, 703, 443), fill="#0b7285", width=5)

    draw_round_rect(draw, (900, 270, 1530, 800), 22, "#f8fbfe", "#dce8f0", 2)
    draw.text((930, 295), "Centres recommandes", font=FONT_H2, fill="#0f172a")
    cards = [
        ("Hopital Regional A", "2.4 km", "Urgences, imagerie, laboratoire"),
        ("Centre Medical B", "4.1 km", "Consultation generale, vaccination"),
        ("Maternite C", "6.8 km", "Maternite, pediatrie, echographie"),
    ]
    cy = 350
    for name, km, details in cards:
        draw_round_rect(draw, (930, cy, 1490, cy + 120), 16, "white", "#cfdce6", 2)
        draw.text((955, cy + 18), f"{name} - {km}", font=FONT_H3, fill="#0f172a")
        draw_multiline(draw, (955, cy + 52), details, FONT_SMALL, "#475569", 500)
        cy += 145
    path = SHOTS_DIR / "02-nearby-centers.png"
    img.save(path)
    return path


def save_emergency():
    img, draw = new_canvas("Coordination des urgences", "SAMU, pompiers et services de reponse")
    filters = ["Non traite", "En cours", "Rejete", "Traite", "Tous"]
    fx = 60
    for item in filters:
        active = item == "En cours"
        draw_round_rect(draw, (fx, 150, fx + 155, 195), 14, "#0c5a7a" if active else "#ffffff", "#cfdce6", 2)
        draw.text((fx + 20, 163), item, font=FONT_SMALL, fill="white" if active else "#1e293b")
        fx += 175
    draw.text((1000, 160), "Periode: 01/04/2026 - 23/04/2026", font=FONT_SMALL, fill="#475569")

    alerts = [
        ("Accident de circulation", "EN ROUTE", "#fff7ed", "#c2410c"),
        ("Evacuation obstetricale", "ON SITE", "#eff6ff", "#1d4ed8"),
        ("Detresse respiratoire", "NEW", "#fef2f2", "#b91c1c"),
    ]
    y = 240
    for title, status, bg, fg in alerts:
        draw_round_rect(draw, (60, y, 990, y + 170), 18, bg, "#d8e4ef", 2)
        draw.text((90, y + 24), title, font=FONT_H3, fill="#0f172a")
        draw_round_rect(draw, (730, y + 20, 900, y + 56), 14, "white", fg, 2)
        draw.text((760, y + 30), status, font=FONT_SMALL, fill=fg)
        draw.text((90, y + 70), "Service cible: SAMU", font=FONT_BODY, fill="#334155")
        draw.text((90, y + 98), "Telephone: +225 07 00 00 00 00", font=FONT_BODY, fill="#334155")
        draw.text((90, y + 126), "Position: 5.37, -4.01 | Point de prise en charge: Centre de sante", font=FONT_SMALL, fill="#475569")
        y += 190

    draw_round_rect(draw, (1040, 240, 1530, 800), 18, "#f8fbfe", "#d8e4ef", 2)
    draw.text((1070, 270), "Mini-carte et suivi", font=FONT_H2, fill="#0f172a")
    draw.rectangle((1080, 330, 1490, 560), fill="#e2f1e8")
    draw.line((1110, 350, 1460, 540), fill="#c0d8c6", width=3)
    draw.line((1450, 360, 1120, 530), fill="#c0d8c6", width=3)
    draw.ellipse((1240, 420, 1272, 452), fill="#dc2626", outline="white", width=3)
    draw.ellipse((1340, 480, 1370, 510), fill="#0ea5e9", outline="white", width=3)
    draw.line((1256, 436, 1355, 495), fill="#0b7285", width=4)
    draw_multiline(
        draw,
        (1080, 600),
        "Le module urgence suit les statuts de prise en charge, la position de l'alerte, "
        "les notes d'avancement et la navigation vers le lieu d'intervention.",
        FONT_BODY,
        "#334155",
        390,
    )
    path = SHOTS_DIR / "03-emergency-alerts.png"
    img.save(path)
    return path


def save_complaints():
    img, draw = new_canvas("Plaintes et satisfaction", "Redevabilite et amelioration continue")
    stats = [
        ("Niveau Region", "Portee"),
        ("128", "Centres du perimetre"),
        ("4.3 / 5", "Note moyenne"),
        ("87%", "Taux de satisfaction"),
    ]
    x = 60
    for value, label in stats:
        draw_round_rect(draw, (x, 150, x + 340, 265), 20, "#f8fbfe", "#d8e4ef", 2)
        draw.text((x + 24, 180), value, font=get_font(30, True), fill="#0c5a7a")
        draw.text((x + 24, 225), label, font=FONT_BODY, fill="#475569")
        x += 365

    draw_round_rect(draw, (60, 310, 1540, 390), 16, "#f8fbfe", "#d8e4ef", 2)
    draw_round_rect(draw, (88, 334, 250, 366), 10, "#ffffff", "#cfdce6", 2)
    draw.text((105, 342), "Toutes", font=FONT_SMALL, fill="#334155")
    draw_round_rect(draw, (1270, 332, 1410, 368), 10, "#ffffff", "#cfdce6", 2)
    draw.text((1295, 341), "Actualiser", font=FONT_SMALL, fill="#334155")
    draw_round_rect(draw, (1420, 332, 1510, 368), 10, "#dbeafe", "#93c5fd", 2)
    draw.text((1434, 341), "Synthese", font=FONT_SMALL, fill="#1d4ed8")

    complaints = [
        ("Attente trop longue aux urgences", "IN_PROGRESS", "#fff7ed"),
        ("Accueil satisfaisant a la maternite", "RESOLVED", "#edfdf4"),
        ("Rupture temporaire de medicaments", "NEW", "#eff6ff"),
    ]
    y = 430
    for title, status, bg in complaints:
        draw_round_rect(draw, (60, y, 1540, y + 110), 16, bg, "#d8e4ef", 2)
        draw.text((88, y + 18), title, font=FONT_H3, fill="#0f172a")
        draw.text((88, y + 52), "Centre: Hopital General | Usager: Anonyme", font=FONT_BODY, fill="#334155")
        draw.text((88, y + 78), "Le flux de plaintes permet de tracer les actions et de renforcer la confiance.", font=FONT_SMALL, fill="#64748b")
        draw_round_rect(draw, (1330, y + 20, 1485, y + 54), 12, "white", "#94a3b8", 2)
        draw.text((1360, y + 29), status, font=FONT_SMALL, fill="#334155")
        y += 130
    path = SHOTS_DIR / "04-complaints.png"
    img.save(path)
    return path


def save_settings():
    img, draw = new_canvas("Administration et gouvernance", "Utilisateurs, centres, regions et districts")
    kpis = [("Utilisateurs", "845"), ("Chefs en attente", "14"), ("Comptes desactives", "9"), ("Centres en attente", "37")]
    x = 60
    for label, value in kpis:
        draw_round_rect(draw, (x, 150, x + 340, 255), 18, "#f8fbfe", "#d8e4ef", 2)
        draw.text((x + 24, 175), label, font=FONT_BODY, fill="#475569")
        draw.text((x + 24, 208), value, font=get_font(28, True), fill="#0c5a7a")
        x += 365

    draw_round_rect(draw, (60, 295, 520, 760), 18, "#f8fbfe", "#d8e4ef", 2)
    draw.text((86, 322), "Creation / modification utilisateur", font=FONT_H3, fill="#0f172a")
    fields = ["Nom complet", "Email", "Mot de passe", "Role principal", "Region", "District", "Centre"]
    fy = 370
    for f in fields:
        draw.text((90, fy - 22), f, font=FONT_SMALL, fill="#64748b")
        draw_round_rect(draw, (90, fy, 480, fy + 34), 8, "white", "#cfdce6", 2)
        fy += 62
    draw_round_rect(draw, (90, 705, 220, 742), 10, "#0c5a7a", "#0c5a7a", 2)
    draw.text((130, 715), "Ajouter", font=FONT_SMALL, fill="white")

    draw_round_rect(draw, (560, 295, 1540, 760), 18, "#f8fbfe", "#d8e4ef", 2)
    draw.text((588, 322), "Table de gestion des comptes", font=FONT_H3, fill="#0f172a")
    cols = ["Nom", "Email", "Role", "Affectation", "Statut"]
    cx = [590, 780, 980, 1130, 1350]
    for i, c in enumerate(cols):
        draw.text((cx[i], 370), c, font=FONT_SMALL, fill="#64748b")
    rows = [
        ("A. Kone", "a.kone@gouv.ci", "REGION", "ABJ", "ACTIF"),
        ("M. Kouassi", "m.kouassi@chu.ci", "ETABLISSEMENT", "CHU", "PENDING"),
        ("S. Traore", "s.traore@samu.ci", "SAMU", "DISTRICT 1", "ACTIF"),
        ("D. Nguessan", "d.ng@gouv.ci", "DISTRICT", "YAMOUSSOUKRO", "ACTIF"),
    ]
    ry = 410
    for row in rows:
        draw_round_rect(draw, (585, ry, 1510, ry + 58), 10, "white", "#d8e4ef", 2)
        vals = list(row)
        for i, val in enumerate(vals):
            draw.text((cx[i], ry + 18), val, font=FONT_SMALL, fill="#334155")
        ry += 72
    path = SHOTS_DIR / "05-settings-admin.png"
    img.save(path)
    return path


def build_images():
    SHOTS_DIR.mkdir(parents=True, exist_ok=True)
    return [
        save_mobile_nearby(),
        save_mobile_emergency(),
        save_mobile_complaints(),
        save_mobile_chef(),
        save_overview(),
    ]


def set_run_font(run, size=20, bold=False, color="0F172A"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = __import__("pptx.dml.color").dml.color.RGBColor.from_string(color)


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = __import__("pptx.dml.color").dml.color.RGBColor.from_string("F4F7FB")
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(8.5), Inches(2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run_font(r, 28, True, "0C5A7A")
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = subtitle
    set_run_font(r, 16, False, "475569")
    slide.shapes.add_picture(str(ROOT / "web-vue" / "public" / "logo-512.png"), Inches(10.2), Inches(0.8), width=Inches(2.1))
    return slide


def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(11.8), Inches(0.8))
    r = title_box.text_frame.paragraphs[0].add_run()
    r.text = title
    set_run_font(r, 24, True, "0C5A7A")
    body = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11.2), Inches(5.5))
    tf = body.text_frame
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(12)
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            set_run_font(run, 18, False, "1E293B")
    return slide


def add_image_slide(prs, title, subtitle, image_path, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(11.5), Inches(0.8))
    r = title_box.text_frame.paragraphs[0].add_run()
    r.text = title
    set_run_font(r, 24, True, "0C5A7A")
    sub = slide.shapes.add_textbox(Inches(0.65), Inches(0.95), Inches(11), Inches(0.4))
    r = sub.text_frame.paragraphs[0].add_run()
    r.text = subtitle
    set_run_font(r, 13, False, "64748B")
    slide.shapes.add_picture(str(image_path), Inches(0.7), Inches(1.45), width=Inches(7.4))
    box = slide.shapes.add_textbox(Inches(8.35), Inches(1.65), Inches(4.4), Inches(4.8))
    tf = box.text_frame
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.space_after = Pt(12)
        for run in p.runs:
            set_run_font(run, 17, False, "1E293B")
    return slide


def build_pptx(images):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "SANTE APROXIMITE",
        "Presentation orientee mobile pour autorites sanitaires, gouvernement et bailleurs",
    )
    add_bullets_slide(
        prs,
        "Pourquoi l'approche mobile est decisive",
        [
            "Le smartphone est le point de contact le plus direct entre le citoyen, le centre de sante et l'urgence.",
            "Une experience mobile simple accelere l'adoption sur le terrain et rend l'impact visible tres vite.",
            "Les usages prioritaires sont concrets: chercher un centre, signaler une urgence, deposer une plainte, suivre une reponse.",
            "Cette approche parle fortement aux bailleurs car elle combine utilite sociale, preuves d'usage et capacite de passage a l'echelle.",
        ],
    )
    add_image_slide(
        prs,
        "Vue mobile 1 - Trouver un centre",
        "Experience smartphone de recherche, orientation et navigation",
        images[0],
        [
            "Recherche de centres a proximite par GPS.",
            "Consultation des services, du plateau technique et des avis.",
            "Navigation directe vers le point de prise en charge.",
        ],
    )
    add_image_slide(
        prs,
        "Vue mobile 2 - Alerte d'urgence",
        "Signalement mobile avec GPS, service cible et suivi du statut",
        images[1],
        [
            "Le citoyen envoie une alerte au bon service en quelques secondes.",
            "Le GPS et les informations du terrain reduisent les pertes de temps.",
            "Le suivi de l'alerte rassure l'usager et valorise la reponse publique.",
        ],
    )
    add_image_slide(
        prs,
        "Vue mobile 3 - Plaintes et satisfaction",
        "Participation citoyenne et boucle de redevabilite",
        images[2],
        [
            "Deposer une plainte devient simple et accessible.",
            "L'usager peut suivre l'etat d'avancement et valider la resolution.",
            "Les autorites obtiennent une mesure continue de la qualite percue.",
        ],
    )
    add_image_slide(
        prs,
        "Vue mobile 4 - Responsable de centre",
        "Le chef d'etablissement agit aussi depuis son telephone",
        images[3],
        [
            "Creation et mise a jour du centre depuis mobile.",
            "Suivi des indicateurs de satisfaction et des plaintes.",
            "Reponse rapide aux remontées du terrain.",
        ],
    )
    add_image_slide(
        prs,
        "Vue complementaire - Pilotage web",
        "Le mobile est complete par un back-office de supervision",
        images[4],
        [
            "Les donnees mobiles alimentent les tableaux de bord institutionnels.",
            "Le web sert a la supervision, a la regulation et au reporting.",
            "Le couple mobile + web est la force structurelle de la plateforme.",
        ],
    )
    add_bullets_slide(
        prs,
        "Ce que finance un partenaire en soutenant l'approche mobile",
        [
            "Un acces plus rapide aux soins pour les populations.",
            "Une meilleure remontee des urgences et des plaintes depuis le terrain.",
            "Une infrastructure digitale visible, tangible et facilement demonstrable.",
            "Une base de donnees utile a la supervision ministerielle et au reporting bailleurs.",
        ],
    )
    add_bullets_slide(
        prs,
        "Prochaine etape recommandee",
        [
            "Organiser une demonstration institutionnelle guidee.",
            "Choisir un territoire pilote et definir les indicateurs de succes.",
            "Structurer un partenariat en trois volets: technologie, deploiement, accompagnement.",
            "Lancer ensuite la communication officielle et la mobilisation des bailleurs.",
        ],
    )
    prs.save(PPTX_PATH)


def add_doc_heading(doc, text, level=1):
    doc.add_heading(text, level=level)


def add_doc_bullets(doc, bullets):
    for item in bullets:
        doc.add_paragraph(item, style="List Bullet")


def build_docx(images):
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = DocxPt(11)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("SANTE APROXIMITE")
    r.bold = True
    r.font.size = DocxPt(22)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run("Dossier de presentation et de publicite institutionnelle illustre").italic = True

    add_doc_heading(doc, "Resume executif", 1)
    doc.add_paragraph(
        "Sante Aproximite est une solution de sante digitale dont la force immediate "
        "est l'experience mobile. Depuis un smartphone, un usager peut trouver un centre, "
        "signaler une urgence, deposer une plainte et suivre les reponses."
    )
    add_doc_bullets(
        doc,
        [
            "Recherche mobile des centres de sante a proximite.",
            "Signalement d'urgence avec telephone, GPS, description et suivi.",
            "Plaintes usagers et retour de satisfaction depuis le mobile.",
            "Appui au responsable de centre pour mettre a jour son etablissement.",
        ],
    )

    sections = [
        ("Vue mobile 1 - Trouver un centre et lancer l'itineraire", images[0]),
        ("Vue mobile 2 - Signaler une urgence et suivre la prise en charge", images[1]),
        ("Vue mobile 3 - Deposer une plainte et suivre son traitement", images[2]),
        ("Vue mobile 4 - Gerer son centre depuis le smartphone", images[3]),
        ("Vue complementaire - Pilotage et supervision web", images[4]),
    ]
    for title, path in sections:
        add_doc_heading(doc, title, 1)
        doc.add_picture(str(path), width=DocxInches(6.5))
        doc.add_paragraph(
            "Illustration representative de l'interface actuelle de l'application, "
            "mise en forme pour un usage de presentation institutionnelle."
        )

    add_doc_heading(doc, "Pourquoi une presentation mobile convainc davantage", 1)
    add_doc_bullets(
        doc,
        [
            "Le smartphone met le service public au plus pres du citoyen.",
            "Les usages mobiles sont plus faciles a demonstrer en reunion institutionnelle.",
            "L'impact terrain devient immediatement compréhensible pour un bailleur.",
            "Le mobile favorise l'adoption rapide dans les zones ou l'ordinateur est moins accessible.",
        ],
    )
    add_doc_heading(doc, "Proposition de partenariat", 1)
    add_doc_bullets(
        doc,
        [
            "Pilote territorial sur 3 a 6 mois.",
            "Financement du passage a l'echelle.",
            "Appui equipements, connectivite et formation.",
            "Co-construction des tableaux de bord et du reporting d'impact.",
        ],
    )
    doc.save(DOCX_PATH)


def main():
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    images = build_images()
    build_pptx(images)
    build_docx(images)
    print(f"PPT cree: {PPTX_PATH}")
    print(f"DOCX cree: {DOCX_PATH}")


if __name__ == "__main__":
    main()
