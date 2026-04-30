"""
Generates a premium commercial pitch deck for Santé Aproximite.
Uses real screenshots from docs/mobile-real-shots-expo/ and docs/presentation-assets/.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.oxml.ns import qn
from lxml import etree
from copy import deepcopy

# ── Paths ──────────────────────────────────────────────────────────────────────
BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
EXPO = os.path.join(BASE, "docs", "mobile-real-shots-expo")
ASSETS = os.path.join(BASE, "docs", "presentation-assets")
OUT = os.path.join(BASE, "docs", "Pitch_Deck_Sante_Aproximite_2025.pptx")

# ── Brand colours ──────────────────────────────────────────────────────────────
DARK_NAVY   = RGBColor(0x1C, 0x23, 0x40)   # app header
TEAL        = RGBColor(0x00, 0x96, 0x88)   # primary teal
ORANGE      = RGBColor(0xE8, 0x62, 0x0A)   # CTA orange
PURPLE      = RGBColor(0x7B, 0x5E, 0xA7)   # security purple
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF5, 0xF6, 0xFA)
MID_GRAY    = RGBColor(0x8E, 0x99, 0xB0)
ACCENT_GREEN= RGBColor(0x2E, 0xCC, 0x71)

# ── Slide dimensions (16:9) ────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ──────────────────────────────────────────────────────────────────────────────
# Helper utilities
# ──────────────────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb=None, fill_alpha=None, line_rgb=None, line_width=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)          # MSO_SHAPE_TYPE.RECTANGLE
    fill = shape.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = shape.line
    if line_rgb:
        line.color.rgb = line_rgb
        line.width = line_width or Pt(1)
    else:
        line.fill.background()
    return shape


def add_textbox(slide, text, x, y, w, h,
                font_name="Segoe UI", font_size=18, bold=False, italic=False,
                color=WHITE, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_image(slide, img_path, x, y, w=None, h=None):
    if not os.path.exists(img_path):
        print(f"  [WARN] image not found: {img_path}")
        return None
    if w and h:
        return slide.shapes.add_picture(img_path, x, y, width=w, height=h)
    elif w:
        return slide.shapes.add_picture(img_path, x, y, width=w)
    elif h:
        return slide.shapes.add_picture(img_path, x, y, height=h)
    else:
        return slide.shapes.add_picture(img_path, x, y)


def add_circle(slide, cx, cy, r, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(9, cx - r, cy - r, 2*r, 2*r)  # OVAL
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape


def gradient_rect(slide, x, y, w, h, color1, color2):
    """Approximate gradient with two overlapping semi-transparent rects."""
    r1 = add_rect(slide, x, y, w, h, fill_rgb=color1)
    return r1


def phone_frame(slide, img_path, cx, cy, phone_h=Inches(5.6)):
    """Draw a phone silhouette and place the screenshot inside it."""
    ratio = 9/19.5
    phone_w = phone_h * ratio
    px = cx - phone_w / 2
    py = cy - phone_h / 2

    # shadow
    shadow = slide.shapes.add_shape(1, px + Inches(0.08), py + Inches(0.08), phone_w, phone_h)
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(0x00,0x00,0x00)
    shadow.line.fill.background()
    # crude transparency via XML
    sp_elem = shadow.fill._xPr
    # we'll just use dark navy opacity workaround – keep it simple

    # phone body
    body = slide.shapes.add_shape(5, px, py, phone_w, phone_h)  # ROUNDED_RECT
    body.adjustments[0] = 0.08
    body.fill.solid()
    body.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x2A)
    body.line.color.rgb = RGBColor(0x44, 0x44, 0x55)
    body.line.width = Pt(2)

    # screen inset
    pad_x = phone_w * 0.06
    pad_y = phone_h * 0.07
    sx = px + pad_x
    sy = py + pad_y
    sw = phone_w - 2*pad_x
    sh = phone_h - 2*pad_y

    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path, sx, sy, width=sw, height=sh)

    # home bar
    bar_w = phone_w * 0.30
    bar = slide.shapes.add_shape(1, cx - bar_w/2, py + phone_h - Inches(0.18), bar_w, Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x88,0x88,0x99)
    bar.line.fill.background()


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 1 – Cover
# ──────────────────────────────────────────────────────────────────────────────
def slide_cover():
    sl = prs.slides.add_slide(BLANK)

    # Full dark navy background
    add_rect(sl, 0, 0, W, H, DARK_NAVY)

    # Decorative teal arc (large circle off-screen top-left)
    arc = sl.shapes.add_shape(9, Inches(-2), Inches(-2), Inches(7), Inches(7))
    arc.fill.solid()
    arc.fill.fore_color.rgb = RGBColor(0x00, 0x88, 0x7A)
    arc.line.fill.background()

    # Decorative orange circle bottom-right
    arc2 = sl.shapes.add_shape(9, Inches(9.5), Inches(4.5), Inches(5), Inches(5))
    arc2.fill.solid()
    arc2.fill.fore_color.rgb = RGBColor(0xB0, 0x46, 0x00)
    arc2.line.fill.background()

    # White card centre
    card_x, card_y = Inches(0.9), Inches(0.8)
    card_w, card_h = Inches(6.5), Inches(5.9)
    card = add_rect(sl, card_x, card_y, card_w, card_h, WHITE)
    card.adjustments  # no rounded corners needed

    # App icon placeholder (teal square with + pin)
    icon_size = Inches(1.0)
    icon_x = card_x + card_w/2 - icon_size/2
    icon_y = card_y + Inches(0.4)
    icon_bg = sl.shapes.add_shape(5, icon_x, icon_y, icon_size, icon_size)
    icon_bg.adjustments[0] = 0.15
    icon_bg.fill.solid()
    icon_bg.fill.fore_color.rgb = DARK_NAVY
    icon_bg.line.fill.background()
    add_textbox(sl, "📍", icon_x, icon_y + Inches(0.05), icon_size, icon_size,
                font_size=36, color=TEAL, align=PP_ALIGN.CENTER)

    # App name
    add_textbox(sl, "Santé Aproximite", card_x, card_y + Inches(1.6), card_w, Inches(0.9),
                font_size=36, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)

    # Tagline
    add_textbox(sl, "La santé de proximité dans votre poche",
                card_x, card_y + Inches(2.5), card_w, Inches(0.5),
                font_size=16, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Badges
    badges = [("🏥 Urgence 24/7", TEAL), ("📍 Géolocalisation", ORANGE), ("🛡️ Sécurité", PURPLE)]
    bx = card_x + Inches(0.5)
    for label, col in badges:
        b = sl.shapes.add_shape(5, bx, card_y + Inches(3.2), Inches(1.7), Inches(0.38))
        b.adjustments[0] = 0.5
        b.fill.solid()
        b.fill.fore_color.rgb = col
        b.line.fill.background()
        add_textbox(sl, label, bx, card_y + Inches(3.2), Inches(1.7), Inches(0.38),
                    font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        bx += Inches(1.85)

    # Made by
    add_textbox(sl, "Par YEFA TECHNOLOGIE  ·  yefa.technologie@gmail.com",
                card_x, card_y + Inches(5.1), card_w, Inches(0.4),
                font_size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Right side – phone mockup with first screenshot
    phone_cx = Inches(10.3)
    phone_cy = Inches(3.75)
    phone_frame(sl, os.path.join(EXPO, "08-chef-login-result.png"), phone_cx, phone_cy, Inches(5.4))

    # PITCH DECK label top-right
    add_textbox(sl, "PITCH DECK  ·  2025", Inches(9.5), Inches(0.2), Inches(3.5), Inches(0.4),
                font_size=11, color=RGBColor(0x99,0xAA,0xBB), align=PP_ALIGN.RIGHT)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 2 – Le Problème
# ──────────────────────────────────────────────────────────────────────────────
def slide_probleme():
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, LIGHT_GRAY)

    # Left accent bar
    add_rect(sl, 0, 0, Inches(0.12), H, ORANGE)

    # Section tag
    add_textbox(sl, "LE PROBLÈME", Inches(0.3), Inches(0.3), Inches(4), Inches(0.4),
                font_size=11, bold=True, color=ORANGE, align=PP_ALIGN.LEFT)

    # Title
    add_textbox(sl, "L'accès aux soins en Afrique :\nune réalité difficile", Inches(0.3), Inches(0.7),
                Inches(6), Inches(1.4), font_size=30, bold=True, color=DARK_NAVY)

    problems = [
        ("😰", "Trouver un centre de santé équipé", "Impossible de savoir à l'avance les services disponibles ni les distances réelles."),
        ("🚨", "Urgences sans coordination", "Pas de système unifié pour alerter le SAMU, les pompiers ou la police et suivre l'intervention."),
        ("😤", "Plaintes sans réponse", "Les patients n'ont aucun canal officiel pour signaler les mauvaises pratiques."),
        ("📵", "Déconnexion dans les zones rurales", "Les réseaux instables coupent l'accès aux services critiques au pire moment."),
    ]

    cols = [(Inches(0.3), Inches(2.3)), (Inches(4.5), Inches(2.3)),
            (Inches(0.3), Inches(4.5)), (Inches(4.5), Inches(4.5))]

    for (x, y), (icon, title, desc) in zip(cols, problems):
        card = add_rect(sl, x, y, Inches(3.9), Inches(1.9), WHITE)
        add_textbox(sl, icon, x + Inches(0.15), y + Inches(0.15), Inches(0.6), Inches(0.6), font_size=26, color=DARK_NAVY)
        add_textbox(sl, title, x + Inches(0.15), y + Inches(0.75), Inches(3.6), Inches(0.35),
                    font_size=13, bold=True, color=DARK_NAVY)
        add_textbox(sl, desc, x + Inches(0.15), y + Inches(1.1), Inches(3.6), Inches(0.7),
                    font_size=10, color=MID_GRAY)

    # Right – big stat
    stat_x = Inches(8.8)
    add_rect(sl, stat_x, Inches(0.5), Inches(4.2), Inches(6.5), DARK_NAVY)
    add_textbox(sl, "Le saviez-vous ?", stat_x, Inches(0.9), Inches(4.2), Inches(0.5),
                font_size=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    stats = [
        ("60%", "des Africains n'ont pas accès\nà un médecin à moins de 5 km"),
        ("8 min", "de délai moyen gagné quand\nl'urgence est bien géolocalisée"),
        ("3 124", "centres de santé déjà\nrépertoriés dans la base"),
    ]
    sy = Inches(1.7)
    for val, label in stats:
        add_textbox(sl, val, stat_x, sy, Inches(4.2), Inches(0.7),
                    font_size=34, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
        add_textbox(sl, label, stat_x, sy + Inches(0.65), Inches(4.2), Inches(0.6),
                    font_size=11, color=WHITE, align=PP_ALIGN.CENTER)
        sy += Inches(1.7)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 3 – Notre Solution
# ──────────────────────────────────────────────────────────────────────────────
def slide_solution():
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, DARK_NAVY)
    add_rect(sl, 0, 0, Inches(0.12), H, TEAL)

    add_textbox(sl, "NOTRE SOLUTION", Inches(0.3), Inches(0.3), Inches(5), Inches(0.4),
                font_size=11, bold=True, color=TEAL)
    add_textbox(sl, "Une plateforme tout-en-un\npour la santé de proximité",
                Inches(0.3), Inches(0.7), Inches(7.5), Inches(1.4),
                font_size=30, bold=True, color=WHITE)
    add_textbox(sl,
        "Santé Aproximite connecte citoyens, établissements de santé et secours\n"
        "via une application mobile intuitive et hors-ligne, géolocalisée en temps réel.",
        Inches(0.3), Inches(2.1), Inches(7.5), Inches(0.8),
        font_size=14, color=RGBColor(0xBB,0xCC,0xDD))

    pillars = [
        (TEAL,   "🏥", "Découverte",   "Trouvez le centre le plus proche\navec ses services et équipements"),
        (ORANGE, "🚑", "Urgences",     "Signalez et suivez toute urgence\nmédical ou sécuritaire en temps réel"),
        (PURPLE, "📋", "Plaintes",     "Déposez et suivez vos plaintes\navec réponse officielle garantie"),
        (ACCENT_GREEN, "📊", "Gestion","Pilotez votre établissement\nstatistiques, notation, approbation"),
    ]

    px = Inches(0.3)
    for col, icon, title, desc in pillars:
        add_rect(sl, px, Inches(3.2), Inches(3.0), Inches(3.8), RGBColor(0x25,0x2E,0x50))
        # top accent
        add_rect(sl, px, Inches(3.2), Inches(3.0), Inches(0.12), col)
        add_textbox(sl, icon, px + Inches(0.2), Inches(3.4), Inches(0.8), Inches(0.7), font_size=30, color=WHITE)
        add_textbox(sl, title, px + Inches(0.2), Inches(4.15), Inches(2.7), Inches(0.45),
                    font_size=15, bold=True, color=WHITE)
        add_textbox(sl, desc, px + Inches(0.2), Inches(4.6), Inches(2.7), Inches(0.9),
                    font_size=10, color=RGBColor(0xAA,0xBB,0xCC))
        px += Inches(3.22)

    # phone on right
    phone_frame(sl, os.path.join(EXPO, "01-centres.png"), Inches(10.5), Inches(3.9), Inches(5.0))


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 4 – Fonctionnalité : Centres de Santé
# ──────────────────────────────────────────────────────────────────────────────
def slide_centres():
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, LIGHT_GRAY)
    add_rect(sl, 0, 0, W, Inches(1.5), DARK_NAVY)
    add_rect(sl, 0, 0, Inches(0.12), H, TEAL)

    add_textbox(sl, "FONCTIONNALITÉ 1", Inches(0.3), Inches(0.22), Inches(5), Inches(0.35),
                font_size=10, bold=True, color=TEAL)
    add_textbox(sl, "Centres de Santé à Proximité", Inches(0.3), Inches(0.55), Inches(8), Inches(0.75),
                font_size=26, bold=True, color=WHITE)

    points = [
        ("📍", "Carte GPS interactive", "Visualisez les centres autour de vous avec des marqueurs colorés sur Google Maps."),
        ("🔍", "Recherche filtrée", "Filtrez par nom de centre ou service spécifique. Rayon ajustable de 1 à 700 km."),
        ("⭐", "Notes & Satisfaction", "Consultez les avis d'autres patients, notez de 1 à 5 étoiles, signalez votre satisfaction."),
        ("🚗", "Itinéraire en 1 clic", "Lance directement Google Maps pour vous guider jusqu'au centre choisi."),
        ("📡", "Données hors-ligne", "3 124 centres synchronisés localement — consultables même sans réseau."),
        ("🏗️", "Plateau technique", "Connaissez les équipements disponibles avant même de vous déplacer."),
    ]

    py = Inches(1.75)
    for i, (icon, title, desc) in enumerate(points):
        col_x = Inches(0.3) if i % 2 == 0 else Inches(4.3)
        if i % 2 == 0 and i > 0:
            py += Inches(1.55)
        add_rect(sl, col_x, py, Inches(3.7), Inches(1.35), WHITE)
        add_rect(sl, col_x, py, Inches(0.08), Inches(1.35), TEAL)
        add_textbox(sl, icon, col_x + Inches(0.2), py + Inches(0.1), Inches(0.6), Inches(0.5), font_size=22, color=DARK_NAVY)
        add_textbox(sl, title, col_x + Inches(0.85), py + Inches(0.1), Inches(2.7), Inches(0.38),
                    font_size=12, bold=True, color=DARK_NAVY)
        add_textbox(sl, desc, col_x + Inches(0.85), py + Inches(0.48), Inches(2.7), Inches(0.75),
                    font_size=9.5, color=MID_GRAY)

    # Screenshots
    phone_frame(sl, os.path.join(EXPO, "01-centres.png"), Inches(10.2), Inches(4.0), Inches(5.0))


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 5 – Fonctionnalité : Urgences Sanitaires
# ──────────────────────────────────────────────────────────────────────────────
def slide_urgences():
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, LIGHT_GRAY)
    add_rect(sl, 0, 0, W, Inches(1.5), RGBColor(0x8B, 0x00, 0x00))
    add_rect(sl, 0, 0, Inches(0.12), H, ORANGE)

    add_textbox(sl, "FONCTIONNALITÉ 2", Inches(0.3), Inches(0.22), Inches(5), Inches(0.35),
                font_size=10, bold=True, color=ORANGE)
    add_textbox(sl, "Urgences Sanitaires 24/7", Inches(0.3), Inches(0.55), Inches(8), Inches(0.75),
                font_size=26, bold=True, color=WHITE)

    steps = [
        ("01", ORANGE, "Choisissez le service",
         "SAMU (médical) ou Sapeurs-Pompiers (incendie, accidents, secours)"),
        ("02", RGBColor(0xFF,0x80,0x00), "Décrivez l'urgence",
         "Type d'urgence parmi 8 catégories + description libre + photos optionnelles"),
        ("03", RGBColor(0xFF,0xA0,0x00), "Capturez votre position GPS",
         "Précision ≤ 300 m exigée — code couleur pour valider la précision"),
        ("04", ACCENT_GREEN, "Suivi en temps réel",
         "Nouveau → Pris en charge → En route → Sur place → Terminé"),
    ]

    sy = Inches(1.75)
    for num, col, title, desc in steps:
        add_rect(sl, Inches(0.3), sy, Inches(7.6), Inches(1.15), WHITE)
        # number badge
        badge = sl.shapes.add_shape(9, Inches(0.45), sy + Inches(0.25), Inches(0.65), Inches(0.65))
        badge.fill.solid()
        badge.fill.fore_color.rgb = col
        badge.line.fill.background()
        add_textbox(sl, num, Inches(0.45), sy + Inches(0.25), Inches(0.65), Inches(0.65),
                    font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(sl, title, Inches(1.3), sy + Inches(0.1), Inches(6.4), Inches(0.42),
                    font_size=14, bold=True, color=DARK_NAVY)
        add_textbox(sl, desc, Inches(1.3), sy + Inches(0.52), Inches(6.4), Inches(0.55),
                    font_size=11, color=MID_GRAY)
        sy += Inches(1.28)

    # Quick-dial callout
    add_rect(sl, Inches(0.3), Inches(6.65), Inches(7.6), Inches(0.65), RGBColor(0xFF,0xEE,0xDD))
    add_textbox(sl, "📞  Appel rapide intégré :  SAMU 15   |   Pompiers 18",
                Inches(0.3), Inches(6.65), Inches(7.6), Inches(0.65),
                font_size=13, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    phone_frame(sl, os.path.join(EXPO, "04-urgence-sanitaire.png"), Inches(10.2), Inches(4.0), Inches(5.0))


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 6 – Fonctionnalité : Urgences Sécuritaires
# ──────────────────────────────────────────────────────────────────────────────
def slide_securite():
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, LIGHT_GRAY)
    add_rect(sl, 0, 0, W, Inches(1.5), PURPLE)
    add_rect(sl, 0, 0, Inches(0.12), H, PURPLE)

    add_textbox(sl, "FONCTIONNALITÉ 3", Inches(0.3), Inches(0.22), Inches(5), Inches(0.35),
                font_size=10, bold=True, color=RGBColor(0xCC,0xAA,0xFF))
    add_textbox(sl, "Urgences Sécuritaires", Inches(0.3), Inches(0.55), Inches(8), Inches(0.75),
                font_size=26, bold=True, color=WHITE)

    services = [
        ("👮", "Police",           PURPLE,              ["Agression", "Vol/braquage", "Attroupement suspect", "Intrusion", "Drogue/fumoir"]),
        ("🪖", "Gendarmerie",      RGBColor(0x4A,0x4A,0x8A), ["Attaque armée", "Braquage", "Coupeur de route", "Menace terroriste", "Conflit"]),
        ("🧡", "Protection Civile",RGBColor(0xCC,0x55,0x00), ["Incendie", "Inondation", "Effondrement", "Fuite de gaz", "Sauvetage"]),
    ]

    sx = Inches(0.3)
    for icon, name, col, types in services:
        add_rect(sl, sx, Inches(1.7), Inches(2.55), Inches(4.5), WHITE)
        add_rect(sl, sx, Inches(1.7), Inches(2.55), Inches(0.5), col)
        add_textbox(sl, icon, sx + Inches(0.05), Inches(1.78), Inches(0.5), Inches(0.4),
                    font_size=20, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(sl, name, sx + Inches(0.2), Inches(2.3), Inches(2.2), Inches(0.45),
                    font_size=15, bold=True, color=DARK_NAVY)
        ty = Inches(2.85)
        for t in types:
            chip = sl.shapes.add_shape(5, sx + Inches(0.2), ty, Inches(2.1), Inches(0.35))
            chip.adjustments[0] = 0.5
            chip.fill.solid()
            chip.fill.fore_color.rgb = RGBColor(0xF0,0xEE,0xFF)
            chip.line.fill.background()
            add_textbox(sl, t, sx + Inches(0.2), ty, Inches(2.1), Inches(0.35),
                        font_size=9, color=col, align=PP_ALIGN.CENTER)
            ty += Inches(0.48)
        sx += Inches(2.7)

    # Feature callout
    features = ["📍 Position GPS capturée", "📸 Photos jointes", "📱 Suivi de mes alertes", "🔒 Anonymat préservé"]
    fy = Inches(6.4)
    fx = Inches(0.3)
    for f in features:
        add_textbox(sl, f, fx, fy, Inches(1.95), Inches(0.5),
                    font_size=10, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
        fx += Inches(2.0)

    phone_frame(sl, os.path.join(EXPO, "05-urgence-securitaire.png"), Inches(10.2), Inches(4.0), Inches(5.0))


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 7 – Fonctionnalité : Plaintes
# ──────────────────────────────────────────────────────────────────────────────
def slide_plaintes():
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, LIGHT_GRAY)
    add_rect(sl, 0, 0, W, Inches(1.5), RGBColor(0x1A, 0x6B, 0x4A))
    add_rect(sl, 0, 0, Inches(0.12), H, ACCENT_GREEN)

    add_textbox(sl, "FONCTIONNALITÉ 4", Inches(0.3), Inches(0.22), Inches(5), Inches(0.35),
                font_size=10, bold=True, color=ACCENT_GREEN)
    add_textbox(sl, "Gestion Transparente des Plaintes", Inches(0.3), Inches(0.55), Inches(8), Inches(0.75),
                font_size=26, bold=True, color=WHITE)

    # Two-column layout
    # Left: citizen view
    add_textbox(sl, "👤  Côté Citoyen", Inches(0.3), Inches(1.7), Inches(3.8), Inches(0.45),
                font_size=13, bold=True, color=DARK_NAVY)
    citizen_items = [
        "Plainte liée à un centre ou générale",
        "Sujets rapides : Accueil, Attente, Propreté...",
        "Description détaillée avec compteur de caractères",
        "Suivi : NOUVELLE → EN COURS → RÉSOLUE",
        "Feedback de satisfaction post-résolution",
    ]
    cy = Inches(2.2)
    for item in citizen_items:
        add_textbox(sl, f"✓  {item}", Inches(0.3), cy, Inches(3.8), Inches(0.42),
                    font_size=11, color=DARK_NAVY)
        cy += Inches(0.45)

    # Centre divider
    add_rect(sl, Inches(4.25), Inches(1.65), Inches(0.03), Inches(5.4), RGBColor(0xCC,0xDD,0xCC))

    # Right: manager view
    add_textbox(sl, "🏥  Côté Établissement", Inches(4.4), Inches(1.7), Inches(3.8), Inches(0.45),
                font_size=13, bold=True, color=DARK_NAVY)
    mgr_items = [
        "Tableau de bord des plaintes reçues",
        "Indicateurs : note moyenne, taux de satisfaction",
        "Réponse officielle avec explication traçable",
        "Statuts : NOUVELLE / EN COURS / RÉSOLUE / REJETÉE",
        "Statistiques de satisfaction en temps réel",
    ]
    my = Inches(2.2)
    for item in mgr_items:
        add_textbox(sl, f"✓  {item}", Inches(4.4), my, Inches(3.8), Inches(0.42),
                    font_size=11, color=DARK_NAVY)
        my += Inches(0.45)

    # Status badges
    statuses = [("NOUVELLE", RGBColor(0x00,0x80,0xFF)), ("EN COURS", ORANGE),
                ("RÉSOLUE", ACCENT_GREEN), ("REJETÉE", RGBColor(0xCC,0x00,0x00))]
    bx = Inches(0.3)
    for label, col in statuses:
        b = sl.shapes.add_shape(5, bx, Inches(6.55), Inches(1.9), Inches(0.5))
        b.adjustments[0] = 0.5
        b.fill.solid()
        b.fill.fore_color.rgb = col
        b.line.fill.background()
        add_textbox(sl, label, bx, Inches(6.55), Inches(1.9), Inches(0.5),
                    font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        bx += Inches(2.0)

    phone_frame(sl, os.path.join(EXPO, "03-plaintes.png"), Inches(10.2), Inches(4.0), Inches(5.0))


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 8 – Multi-rôles
# ──────────────────────────────────────────────────────────────────────────────
def slide_roles():
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, DARK_NAVY)
    add_rect(sl, 0, 0, Inches(0.12), H, ORANGE)

    add_textbox(sl, "ARCHITECTURE MULTI-RÔLES", Inches(0.3), Inches(0.3), Inches(8), Inches(0.4),
                font_size=11, bold=True, color=ORANGE)
    add_textbox(sl, "Une app, sept profils utilisateur", Inches(0.3), Inches(0.7),
                Inches(9), Inches(0.75), font_size=28, bold=True, color=WHITE)

    roles = [
        ("👤", "Citoyen",         TEAL,   "Consulte, note, signale"),
        ("🏥", "Chef Étab.",      ORANGE, "Gère son centre"),
        ("🚑", "SAMU",            RGBColor(0xFF,0x44,0x44), "Gère les urgences méd."),
        ("🚒", "Sapeurs-Pompiers",RGBColor(0xFF,0x77,0x00), "Gère incendies/accidents"),
        ("👮", "Police",          PURPLE, "Gère incidents sécurité"),
        ("🪖", "Gendarmerie",     RGBColor(0x33,0x66,0xCC), "Gère alertes terrain"),
        ("🛡️", "Protection Civile",ACCENT_GREEN,"Gère catastrophes"),
    ]

    cols_per_row = 4
    rx, ry = Inches(0.3), Inches(1.7)
    for i, (icon, name, col, action) in enumerate(roles):
        cx = rx + (i % cols_per_row) * Inches(3.22)
        cy = ry + (i // cols_per_row) * Inches(2.3)
        card = add_rect(sl, cx, cy, Inches(3.0), Inches(2.0), RGBColor(0x25,0x2E,0x50))
        add_rect(sl, cx, cy, Inches(3.0), Inches(0.1), col)
        add_textbox(sl, icon, cx + Inches(0.1), cy + Inches(0.2), Inches(0.8), Inches(0.7), font_size=28, color=WHITE)
        add_textbox(sl, name, cx + Inches(0.1), cy + Inches(0.95), Inches(2.8), Inches(0.45),
                    font_size=14, bold=True, color=WHITE)
        add_textbox(sl, action, cx + Inches(0.1), cy + Inches(1.4), Inches(2.8), Inches(0.45),
                    font_size=10, color=RGBColor(0xAA,0xBB,0xCC))

    # Screenshot menu
    phone_frame(sl, os.path.join(EXPO, "02-menu.png"), Inches(10.5), Inches(3.8), Inches(4.6))


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 9 – Gestion Établissement (Chef Screen)
# ──────────────────────────────────────────────────────────────────────────────
def slide_chef():
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, LIGHT_GRAY)
    add_rect(sl, 0, 0, W, Inches(1.5), RGBColor(0x8B, 0x4B, 0x00))
    add_rect(sl, 0, 0, Inches(0.12), H, ORANGE)

    add_textbox(sl, "TABLEAU DE BORD ÉTABLISSEMENT", Inches(0.3), Inches(0.22), Inches(8), Inches(0.35),
                font_size=10, bold=True, color=ORANGE)
    add_textbox(sl, "Piloter son établissement en temps réel", Inches(0.3), Inches(0.55),
                Inches(8), Inches(0.75), font_size=26, bold=True, color=WHITE)

    features = [
        ("🏥", "Créer & gérer son centre",
         "Saisie complète : nom, adresse, region, district, niveau, type (CHU, Clinique, ESPC…), plateau technique et services."),
        ("📍", "Coordonnées GPS précises",
         "Saisie manuelle ou capture automatique depuis l'appareil. Indispensable pour apparaître dans la carte des utilisateurs."),
        ("📊", "Statistiques en direct",
         "Note moyenne ★, taux de satisfaction %, nombre total de plaintes reçues — actualisé toutes les 30 minutes."),
        ("📋", "Répondre aux plaintes",
         "Consultez chaque plainte et apportez une explication officielle tracée et horodatée."),
        ("✅", "Processus d'approbation",
         "Statut visible : EN ATTENTE → APPROUVÉ ou REJETÉ par l'autorité régulatrice."),
    ]

    fy = Inches(1.7)
    for icon, title, desc in features:
        add_rect(sl, Inches(0.3), fy, Inches(7.6), Inches(0.98), WHITE)
        add_rect(sl, Inches(0.3), fy, Inches(0.08), Inches(0.98), ORANGE)
        add_textbox(sl, icon, Inches(0.5), fy + Inches(0.08), Inches(0.6), Inches(0.5), font_size=22, color=DARK_NAVY)
        add_textbox(sl, title, Inches(1.2), fy + Inches(0.08), Inches(6.5), Inches(0.38),
                    font_size=13, bold=True, color=DARK_NAVY)
        add_textbox(sl, desc, Inches(1.2), fy + Inches(0.46), Inches(6.5), Inches(0.45),
                    font_size=10, color=MID_GRAY)
        fy += Inches(1.1)

    phone_frame(sl, os.path.join(EXPO, "09-chef-filled.png"), Inches(10.2), Inches(4.0), Inches(5.0))


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 10 – Stack Technologique
# ──────────────────────────────────────────────────────────────────────────────
def slide_tech():
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, DARK_NAVY)
    add_rect(sl, 0, 0, Inches(0.12), H, TEAL)

    add_textbox(sl, "STACK TECHNOLOGIQUE", Inches(0.3), Inches(0.3), Inches(8), Inches(0.4),
                font_size=11, bold=True, color=TEAL)
    add_textbox(sl, "Production-ready · Scalable · Sécurisé",
                Inches(0.3), Inches(0.7), Inches(9), Inches(0.65),
                font_size=26, bold=True, color=WHITE)

    layers = [
        ("📱 Mobile", ORANGE, [
            "React Native + Expo",
            "expo-location (GPS)",
            "expo-image-picker (Photos)",
            "AsyncStorage (Offline)",
            "react-native-maps",
        ]),
        ("🖥️ Web Admin", TEAL, [
            "Vue.js 3 + Vite",
            "Dashboard temps réel",
            "Gestion des centres",
            "Approbation & modération",
        ]),
        ("⚙️ Backend API", PURPLE, [
            "Node.js + Express",
            "PostgreSQL",
            "JWT Auth (7j / 30j refresh)",
            "Haversine géolocalisation",
            "REST API documentée",
        ]),
        ("🌐 Infrastructure", ACCENT_GREEN, [
            "Serveur dédié Linux",
            "IP publique fixe",
            "API : 193.168.173.181:5000",
            "Déploiement continu prêt",
        ]),
    ]

    lx = Inches(0.3)
    for title, col, items in layers:
        add_rect(sl, lx, Inches(1.55), Inches(3.0), Inches(5.6), RGBColor(0x1E,0x28,0x48))
        add_rect(sl, lx, Inches(1.55), Inches(3.0), Inches(0.12), col)
        add_textbox(sl, title, lx + Inches(0.15), Inches(1.75), Inches(2.8), Inches(0.5),
                    font_size=14, bold=True, color=WHITE)
        iy = Inches(2.35)
        for item in items:
            add_textbox(sl, f"▸  {item}", lx + Inches(0.15), iy, Inches(2.7), Inches(0.42),
                        font_size=10.5, color=RGBColor(0xBB,0xCC,0xDD))
            iy += Inches(0.5)
        lx += Inches(3.22)

    # Key differentiators
    add_textbox(sl, "🔑  Points différenciants",
                Inches(0.3), Inches(7.0), Inches(13), Inches(0.35),
                font_size=12, bold=True, color=TEAL)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 11 – Opportunité de Marché
# ──────────────────────────────────────────────────────────────────────────────
def slide_marche():
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, LIGHT_GRAY)
    add_rect(sl, 0, 0, W, Inches(1.5), DARK_NAVY)
    add_rect(sl, 0, 0, Inches(0.12), H, TEAL)

    add_textbox(sl, "OPPORTUNITÉ DE MARCHÉ", Inches(0.3), Inches(0.22), Inches(8), Inches(0.35),
                font_size=10, bold=True, color=TEAL)
    add_textbox(sl, "Un marché massif et sous-digitalisé",
                Inches(0.3), Inches(0.55), Inches(9), Inches(0.75),
                font_size=26, bold=True, color=WHITE)

    metrics = [
        ("1.4 Md", "Population ciblée\n(Afrique subsaharienne)", TEAL),
        ("+75%", "Pénétration smartphone\nen 2025 (GSMA)", ORANGE),
        ("3 124", "Centres déjà\nrépertoriés", ACCENT_GREEN),
        ("$45 Md", "Marché HealthTech\nAfrique d'ici 2030", PURPLE),
    ]

    mx = Inches(0.3)
    for val, label, col in metrics:
        add_rect(sl, mx, Inches(1.7), Inches(3.0), Inches(2.3), WHITE)
        add_rect(sl, mx, Inches(1.7), Inches(3.0), Inches(0.12), col)
        add_textbox(sl, val, mx, Inches(2.05), Inches(3.0), Inches(0.85),
                    font_size=38, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_textbox(sl, label, mx, Inches(2.9), Inches(3.0), Inches(0.65),
                    font_size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)
        mx += Inches(3.22)

    # Target markets
    add_textbox(sl, "Cibles prioritaires de déploiement",
                Inches(0.3), Inches(4.2), Inches(8), Inches(0.45),
                font_size=14, bold=True, color=DARK_NAVY)

    targets = [
        ("🏛️", "Gouvernements & Ministères de la Santé", "Digitaliser le répertoire national des établissements"),
        ("🏥", "Établissements Privés", "Visibilité, réputation et gestion des retours patients"),
        ("🚑", "Services de Secours", "Coordination des interventions d'urgence"),
        ("📱", "Grand Public", "Application gratuite, accès sans compte pour l'essentiel"),
    ]

    ty = Inches(4.8)
    tx = Inches(0.3)
    for icon, title, desc in targets:
        add_rect(sl, tx, ty, Inches(6.1), Inches(0.82), WHITE)
        add_textbox(sl, icon, tx + Inches(0.1), ty + Inches(0.1), Inches(0.5), Inches(0.5), font_size=20, color=DARK_NAVY)
        add_textbox(sl, title, tx + Inches(0.75), ty + Inches(0.06), Inches(5.2), Inches(0.35),
                    font_size=12, bold=True, color=DARK_NAVY)
        add_textbox(sl, desc, tx + Inches(0.75), ty + Inches(0.42), Inches(5.2), Inches(0.32),
                    font_size=10, color=MID_GRAY)
        ty += Inches(0.93)

    # Map placeholder
    add_rect(sl, Inches(6.6), Inches(4.2), Inches(6.5), Inches(3.1), RGBColor(0xEE,0xF5,0xFF))
    add_textbox(sl, "🌍\n\nCôte d'Ivoire · Sénégal · Mali\nBurkina Faso · Cameroun · Bénin\n\n→ Extension régionale planifiée",
                Inches(6.6), Inches(4.2), Inches(6.5), Inches(3.1),
                font_size=13, color=DARK_NAVY, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 12 – Modèle Économique
# ──────────────────────────────────────────────────────────────────────────────
def slide_business():
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, DARK_NAVY)
    add_rect(sl, 0, 0, Inches(0.12), H, ORANGE)

    add_textbox(sl, "MODÈLE ÉCONOMIQUE", Inches(0.3), Inches(0.3), Inches(8), Inches(0.4),
                font_size=11, bold=True, color=ORANGE)
    add_textbox(sl, "Plusieurs leviers de monétisation", Inches(0.3), Inches(0.7),
                Inches(9), Inches(0.65), font_size=26, bold=True, color=WHITE)

    streams = [
        (TEAL,   "🆓", "Freemium Citoyen",
         "Application gratuite pour les utilisateurs finaux.\nPrime pour fonctionnalités avancées (historique étendu, alertes push premium)."),
        (ORANGE, "💼", "SaaS Établissement",
         "Abonnement mensuel pour les chefs d'établissement :\nvisibilité boostée, statistiques détaillées, support prioritaire."),
        (PURPLE, "🏛️", "Licences Gouvernementales",
         "Contrats B2G : Ministères de la Santé, mairies, districts.\nBase de données nationale des établissements, tableaux de bord régionaux."),
        (ACCENT_GREEN, "🤝", "Partenariats",
         "Compagnies d'assurance santé, ONG, agences de développement.\nIntégration API pour systèmes existants."),
    ]

    sx = Inches(0.3)
    for col, icon, title, desc in streams:
        add_rect(sl, sx, Inches(1.55), Inches(3.0), Inches(4.0), RGBColor(0x1E,0x28,0x48))
        add_rect(sl, sx, Inches(1.55), Inches(3.0), Inches(0.12), col)
        add_textbox(sl, icon, sx + Inches(0.15), Inches(1.85), Inches(0.7), Inches(0.65), font_size=28, color=WHITE)
        add_textbox(sl, title, sx + Inches(0.15), Inches(2.6), Inches(2.7), Inches(0.5),
                    font_size=13, bold=True, color=WHITE)
        add_textbox(sl, desc, sx + Inches(0.15), Inches(3.15), Inches(2.7), Inches(1.5),
                    font_size=10, color=RGBColor(0xBB,0xCC,0xDD))
        sx += Inches(3.22)

    # Roadmap
    add_rect(sl, Inches(0.3), Inches(5.8), Inches(12.73), Inches(1.4), RGBColor(0x15,0x1C,0x38))
    add_textbox(sl, "🗓️  Roadmap", Inches(0.5), Inches(5.88), Inches(2), Inches(0.35),
                font_size=12, bold=True, color=TEAL)
    milestones = [
        ("Q2 2025", "MVP lancé · 3 124 centres · Backend production"),
        ("Q3 2025", "Onboarding établissements · Module régulateur"),
        ("Q4 2025", "Partenariat pilote Ministère · Monétisation SaaS"),
        ("2026",    "Expansion sous-régionale · Levée de fonds Série A"),
    ]
    mx = Inches(0.5)
    for period, text in milestones:
        add_textbox(sl, period, mx, Inches(6.3), Inches(1.5), Inches(0.3),
                    font_size=9, bold=True, color=ORANGE)
        add_textbox(sl, text, mx, Inches(6.6), Inches(2.85), Inches(0.45),
                    font_size=9, color=RGBColor(0xBB,0xCC,0xDD))
        mx += Inches(3.22)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 13 – Pourquoi Investir
# ──────────────────────────────────────────────────────────────────────────────
def slide_invest():
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, LIGHT_GRAY)
    add_rect(sl, 0, 0, W, Inches(1.5), DARK_NAVY)
    add_rect(sl, 0, 0, Inches(0.12), H, ORANGE)

    add_textbox(sl, "POURQUOI INVESTIR", Inches(0.3), Inches(0.22), Inches(8), Inches(0.35),
                font_size=10, bold=True, color=ORANGE)
    add_textbox(sl, "5 raisons de rejoindre l'aventure", Inches(0.3), Inches(0.55),
                Inches(9), Inches(0.75), font_size=26, bold=True, color=WHITE)

    reasons = [
        ("🚀", "Produit déjà fonctionnel",
         "MVP en production avec backend réel, 3 124 centres, 6 modules actifs. Pas un concept — une réalité testable dès aujourd'hui."),
        ("🌍", "Marché sous-digitalisé",
         "L'Afrique subsaharienne a besoin de solutions locales, adaptées au contexte réseau, au multilinguisme et aux usages mobiles-first."),
        ("🔒", "Barrières à l'entrée solides",
         "Données géolocalisées propriétaires, réseau de centres référencés, intégrations services de secours — difficile à répliquer rapidement."),
        ("📈", "Scalabilité prouvée",
         "Architecture API REST découplée. L'ajout de nouvelles régions ou pays se fait par configuration, sans refonte du code."),
        ("🤝", "Équipe focalisée & motivée",
         "YEFA TECHNOLOGIE, studio technique spécialisé en solutions digitales pour l'Afrique avec expertise mobile et backend."),
    ]

    ry = Inches(1.7)
    for i, (icon, title, desc) in enumerate(reasons):
        col_x = Inches(0.3) if i % 2 == 0 else Inches(6.6)
        if i % 2 == 0 and i > 0:
            ry += Inches(1.45)
        add_rect(sl, col_x, ry, Inches(6.0), Inches(1.3), WHITE)
        add_rect(sl, col_x, ry, Inches(0.08), Inches(1.3), ORANGE)
        add_textbox(sl, icon, col_x + Inches(0.2), ry + Inches(0.1), Inches(0.6), Inches(0.6), font_size=24, color=DARK_NAVY)
        add_textbox(sl, title, col_x + Inches(0.95), ry + Inches(0.1), Inches(4.8), Inches(0.38),
                    font_size=13, bold=True, color=DARK_NAVY)
        add_textbox(sl, desc, col_x + Inches(0.95), ry + Inches(0.5), Inches(4.8), Inches(0.7),
                    font_size=10, color=MID_GRAY)

    # Last single reason
    ry += Inches(1.45)
    col_x = Inches(0.3)
    add_rect(sl, col_x, ry, Inches(12.73), Inches(1.0), RGBColor(0xFF,0xF0,0xE0))
    add_rect(sl, col_x, ry, Inches(0.08), Inches(1.0), ORANGE)
    icon, title, desc = reasons[4]
    add_textbox(sl, icon, col_x + Inches(0.2), ry + Inches(0.1), Inches(0.6), Inches(0.6), font_size=24, color=DARK_NAVY)
    add_textbox(sl, title, col_x + Inches(0.95), ry + Inches(0.08), Inches(11.5), Inches(0.38),
                font_size=13, bold=True, color=DARK_NAVY)
    add_textbox(sl, desc, col_x + Inches(0.95), ry + Inches(0.46), Inches(11.5), Inches(0.48),
                font_size=10, color=MID_GRAY)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 14 – Contact / CTA
# ──────────────────────────────────────────────────────────────────────────────
def slide_contact():
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, DARK_NAVY)

    # Teal gradient left half
    add_rect(sl, 0, 0, Inches(6.5), H, RGBColor(0x00,0x7A,0x6E))

    # Decorative circles
    c1 = sl.shapes.add_shape(9, Inches(-1.5), Inches(-1.5), Inches(4), Inches(4))
    c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(0x00,0x60,0x56); c1.line.fill.background()
    c2 = sl.shapes.add_shape(9, Inches(3), Inches(4.5), Inches(3.5), Inches(3.5))
    c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x00,0x55,0x4C); c2.line.fill.background()

    # Left content
    add_textbox(sl, "📍", Inches(0.5), Inches(0.8), Inches(1.5), Inches(1.2), font_size=60, color=WHITE)
    add_textbox(sl, "Santé Aproximite", Inches(0.5), Inches(1.9), Inches(5.7), Inches(0.9),
                font_size=32, bold=True, color=WHITE)
    add_textbox(sl, "La santé de proximité dans votre poche",
                Inches(0.5), Inches(2.8), Inches(5.7), Inches(0.55),
                font_size=16, italic=True, color=RGBColor(0xCC,0xEE,0xEA))

    add_textbox(sl, "Prêt à transformer l'accès à la santé\nen Afrique ?",
                Inches(0.5), Inches(3.55), Inches(5.7), Inches(0.9),
                font_size=20, bold=True, color=WHITE)

    # CTA button
    btn = sl.shapes.add_shape(5, Inches(0.5), Inches(4.7), Inches(2.8), Inches(0.65))
    btn.adjustments[0] = 0.4
    btn.fill.solid(); btn.fill.fore_color.rgb = ORANGE; btn.line.fill.background()
    add_textbox(sl, "Prendre contact →", Inches(0.5), Inches(4.7), Inches(2.8), Inches(0.65),
                font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Right contact info
    add_textbox(sl, "Contactez-nous", Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.6),
                font_size=22, bold=True, color=WHITE)

    contacts = [
        ("📧", "Email", "yefa.technologie@gmail.com"),
        ("🏢", "Entreprise", "YEFA TECHNOLOGIE"),
        ("🌐", "Application", "Santé Aproximite — Mobile & Web"),
        ("📍", "Technologie", "React Native · Node.js · PostgreSQL"),
        ("🚀", "Statut", "MVP en production · Prêt pour investissement"),
    ]

    cy = Inches(2.3)
    for icon, label, value in contacts:
        add_textbox(sl, icon, Inches(7.0), cy, Inches(0.5), Inches(0.5), font_size=18, color=TEAL)
        add_textbox(sl, label, Inches(7.6), cy + Inches(0.03), Inches(1.3), Inches(0.35),
                    font_size=10, bold=True, color=MID_GRAY)
        add_textbox(sl, value, Inches(7.6), cy + Inches(0.38), Inches(5.0), Inches(0.38),
                    font_size=12, color=WHITE)
        cy += Inches(0.9)

    # Powered by tag
    add_textbox(sl, "Powered by YEFA TECHNOLOGIE  ·  yefa.technologie@gmail.com",
                Inches(7.0), Inches(7.0), Inches(6.0), Inches(0.35),
                font_size=9, color=MID_GRAY)

    # Phone mockup – login screen
    phone_frame(sl, os.path.join(EXPO, "08-chef-login-result.png"), Inches(4.5), Inches(5.5), Inches(3.8))


# ──────────────────────────────────────────────────────────────────────────────
# BUILD
# ──────────────────────────────────────────────────────────────────────────────
print("Building slides…")
slide_cover()       ; print("  1/14 Cover ✓")
slide_probleme()    ; print("  2/14 Problème ✓")
slide_solution()    ; print("  3/14 Solution ✓")
slide_centres()     ; print("  4/14 Centres de santé ✓")
slide_urgences()    ; print("  5/14 Urgences sanitaires ✓")
slide_securite()    ; print("  6/14 Urgences sécuritaires ✓")
slide_plaintes()    ; print("  7/14 Plaintes ✓")
slide_roles()       ; print("  8/14 Multi-rôles ✓")
slide_chef()        ; print("  9/14 Gestion établissement ✓")
slide_tech()        ; print(" 10/14 Stack technologique ✓")
slide_marche()      ; print(" 11/14 Marché ✓")
slide_business()    ; print(" 12/14 Modèle économique ✓")
slide_invest()      ; print(" 13/14 Pourquoi investir ✓")
slide_contact()     ; print(" 14/14 Contact / CTA ✓")

prs.save(OUT)
print(f"\n✅  Présentation sauvegardée : {OUT}")
