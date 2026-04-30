from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = ROOT / "docs"
SHOTS_DIR = DOCS_DIR / "mobile-real-shots"
OUTPUT = DOCS_DIR / "Presentation_Mobile_Sante_Aproximite_Comptes_Et_Avantages.pptx"


COLORS = {
    "navy": RGBColor(17, 24, 39),
    "blue": RGBColor(37, 99, 235),
    "teal": RGBColor(13, 148, 136),
    "green": RGBColor(22, 163, 74),
    "purple": RGBColor(124, 58, 237),
    "red": RGBColor(220, 38, 38),
    "amber": RGBColor(217, 119, 6),
    "white": RGBColor(255, 255, 255),
    "bg": RGBColor(244, 247, 251),
    "line": RGBColor(217, 226, 236),
    "text": RGBColor(31, 41, 55),
    "muted": RGBColor(100, 116, 139),
}


def add_bg(slide, prs, color):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["text"]
    return box


def add_card(slide, left, top, width, height, fill=COLORS["white"], line=COLORS["line"]):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.1)
    return shape


def add_header(slide, title, subtitle, accent):
    add_card(slide, Inches(0.45), Inches(0.35), Inches(12.35), Inches(0.95), fill=COLORS["white"], line=COLORS["white"])
    add_card(slide, Inches(0.65), Inches(0.56), Inches(1.85), Inches(0.24), fill=accent, line=accent)
    add_text(slide, Inches(0.66), Inches(0.48), Inches(8.5), Inches(0.45), title, size=25, bold=True, color=COLORS["navy"])
    add_text(slide, Inches(0.66), Inches(0.86), Inches(10.0), Inches(0.28), subtitle, size=11.5, color=COLORS["muted"])


def add_phone_frame(slide, image_path, left, top, width, height):
    shell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shell.fill.solid()
    shell.fill.fore_color.rgb = RGBColor(15, 23, 42)
    shell.line.color.rgb = RGBColor(15, 23, 42)
    shell.line.width = Pt(1.2)

    inner_margin = Inches(0.12)
    screen_left = left + inner_margin
    screen_top = top + Inches(0.22)
    screen_width = width - (inner_margin * 2)
    screen_height = height - Inches(0.38)

    slide.shapes.add_picture(str(image_path), screen_left, screen_top, width=screen_width, height=screen_height)
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left + (width / 2) - Inches(0.5),
        top + Inches(0.06),
        Inches(1.0),
        Inches(0.08),
    ).fill.solid()


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, COLORS["navy"])
    add_card(slide, Inches(0.55), Inches(0.55), Inches(12.1), Inches(6.35), fill=COLORS["white"], line=COLORS["white"])
    add_text(slide, Inches(0.95), Inches(0.95), Inches(5.8), Inches(0.6), "Sante Aproximite", size=28, bold=True, color=COLORS["navy"])
    add_text(slide, Inches(0.95), Inches(1.45), Inches(5.8), Inches(0.95), "Version mobile complete\ncomptes, usages et avantages", size=23, bold=True, color=COLORS["blue"])
    add_text(
        slide,
        Inches(0.95),
        Inches(2.45),
        Inches(4.8),
        Inches(2.0),
        "Une plateforme mobile de terrain pour orienter les usagers, signaler les urgences, remonter les plaintes et outiller les equipes sanitaires et securitaires.",
        size=16,
        color=COLORS["text"],
    )
    add_card(slide, Inches(0.95), Inches(5.0), Inches(2.6), Inches(0.48), fill=COLORS["teal"], line=COLORS["teal"])
    add_text(slide, Inches(1.1), Inches(5.07), Inches(2.2), Inches(0.22), "Captures reelles du mobile", size=12.5, bold=True, color=COLORS["white"])
    add_phone_frame(slide, SHOTS_DIR / "current-screen.png", Inches(7.3), Inches(0.82), Inches(4.4), Inches(5.95))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, COLORS["bg"])
    add_header(slide, "Comptes pris en charge", "Tous les profils utiles a l'ecosysteme mobile", COLORS["blue"])
    roles = [
        "Usager / citoyen",
        "Chef d'etablissement",
        "National",
        "Regulateur",
        "Region",
        "District",
        "SAMU",
        "Sapeurs-pompiers",
        "Police",
        "Gendarmerie",
        "Protection civile",
    ]
    y = 1.65
    x = 0.8
    for index, role in enumerate(roles):
        add_card(slide, Inches(x), Inches(y), Inches(2.35), Inches(0.72))
        add_text(slide, Inches(x + 0.18), Inches(y + 0.19), Inches(2.0), Inches(0.24), role, size=15, bold=True, color=COLORS["navy"], align=PP_ALIGN.CENTER)
        x += 2.55
        if (index + 1) % 4 == 0:
            x = 0.8
            y += 0.95
    add_text(slide, Inches(0.85), Inches(5.65), Inches(11.5), Inches(0.5), "La version mobile ne sert pas seulement l'usager : elle couvre aussi la supervision, la gestion des centres et les equipes d'intervention.", size=15, color=COLORS["text"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, COLORS["bg"])
    add_header(slide, "Parcours d'acces mobile", "Connexion differenciee selon le type de compte", COLORS["amber"])
    add_phone_frame(slide, SHOTS_DIR / "app-launch-phone.png", Inches(0.8), Inches(1.55), Inches(4.0), Inches(5.3))
    add_text(slide, Inches(5.2), Inches(1.9), Inches(6.6), Inches(0.4), "Ce que cet ecran apporte", size=19, bold=True, color=COLORS["navy"])
    bullets = [
        "Choix rapide du profil : utilisateur simple, chef d'etablissement, SAMU / pompiers.",
        "Point d'entree clair pour les usages grand public et professionnels.",
        "Connexion legere par numero pour l'usager et connexion securisee par email pour les comptes professionnels.",
        "Une experience d'accueil simple, lisible et facile a expliquer a grande echelle.",
    ]
    y = 2.45
    for bullet in bullets:
        add_card(slide, Inches(5.25), Inches(y), Inches(0.18), Inches(0.18), fill=COLORS["amber"], line=COLORS["amber"])
        add_text(slide, Inches(5.55), Inches(y - 0.02), Inches(6.0), Inches(0.55), bullet, size=14, color=COLORS["text"])
        y += 0.83

    def image_slide(title, subtitle, image_name, accent, bullets):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide, prs, COLORS["bg"])
        add_header(slide, title, subtitle, accent)
        add_phone_frame(slide, SHOTS_DIR / image_name, Inches(0.8), Inches(1.5), Inches(4.3), Inches(5.45))
        add_text(slide, Inches(5.5), Inches(1.85), Inches(6.1), Inches(0.4), "Valeur mobile", size=19, bold=True, color=COLORS["navy"])
        y = 2.35
        for bullet in bullets:
            add_card(slide, Inches(5.52), Inches(y), Inches(0.18), Inches(0.18), fill=accent, line=accent)
            add_text(slide, Inches(5.82), Inches(y - 0.02), Inches(5.7), Inches(0.6), bullet, size=14, color=COLORS["text"])
            y += 0.82

    image_slide(
        "Acces aux centres de sante",
        "Recherche, proximite, navigation et satisfaction citoyenne",
        "current-screen.png",
        COLORS["teal"],
        [
            "Localisation immediate des centres proches autour de l'usager.",
            "Consultation des services, du plateau technique et de la distance.",
            "Evaluation et satisfaction directement depuis le mobile.",
            "Navigation terrain vers le centre choisi sans passer par un agent.",
        ],
    )

    image_slide(
        "Menu citoyen mobile",
        "Les usages du quotidien disponibles dans une seule application",
        "menu-open.png",
        COLORS["blue"],
        [
            "Centres de sante, plaintes, suivi des plaintes, urgences sanitaires et urgences securitaires.",
            "Navigation simple entre modules sur le terrain.",
            "Une experience qui regroupe orientation, alerte et redevabilite citoyenne.",
            "Le meme mobile devient un guichet sanitaire de proximite.",
        ],
    )

    image_slide(
        "Urgence sanitaire",
        "Signalement cible vers le bon service de secours",
        "urgence-sanitaire.png",
        COLORS["red"],
        [
            "Choix du service de prise en charge : SAMU ou sapeurs-pompiers.",
            "Typologie d'urgence adaptee au service selectionne.",
            "Ajout du lieu, de la description, de photos et de la position GPS.",
            "Reponse plus rapide grace a un signalement mieux qualifie.",
        ],
    )

    image_slide(
        "Urgence securitaire",
        "Alerte terrain pour police, gendarmerie ou protection civile",
        "urgence-securitaire.png",
        COLORS["purple"],
        [
            "Adressage direct a l'autorite competente selon la menace.",
            "Photos, localisation GPS et description de l'incident.",
            "Utilisable pour les agressions, affrontements, troubles ou sinistres.",
            "Un outil utile pour la securite communautaire et la gestion de crise.",
        ],
    )

    image_slide(
        "Plaintes et retour citoyen",
        "Un dispositif mobile de redevabilite et d'amelioration continue",
        "plaintes.png",
        COLORS["blue"],
        [
            "Recherche d'un centre ou plainte generale selon le cas.",
            "Sujets guides pour standardiser les remontes usagers.",
            "Suivi mobile des plaintes et de leur traitement.",
            "Renforcement de la confiance entre population et services de sante.",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, COLORS["bg"])
    add_header(slide, "Avantages strategiques de la version mobile", "Ce que le mobile apporte concretement au systeme", COLORS["green"])
    advantages = [
        ("Accessibilite", "Le service est au plus proche de l'usager, directement dans sa poche."),
        ("Rapidité", "Moins de friction pour signaler, chercher, evaluer ou se faire orienter."),
        ("Terrain", "Les informations captent la realite du terrain : GPS, photos, contexte."),
        ("Synchronisation", "La base locale permet un usage plus robuste et une mise a jour progressive."),
        ("Redevabilite", "Plaintes, satisfaction et notes donnent une voix au citoyen."),
        ("Coordination", "Les urgences sanitaires et securitaires sont mieux routees vers les bons acteurs."),
    ]
    y = 1.75
    for title, body in advantages:
        add_card(slide, Inches(0.8), Inches(y), Inches(11.75), Inches(0.78))
        add_text(slide, Inches(1.05), Inches(y + 0.14), Inches(2.1), Inches(0.25), title, size=16, bold=True, color=COLORS["navy"])
        add_text(slide, Inches(3.0), Inches(y + 0.12), Inches(9.1), Inches(0.32), body, size=13.5, color=COLORS["text"])
        y += 0.9

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, COLORS["navy"])
    add_text(slide, Inches(0.95), Inches(1.2), Inches(10.5), Inches(0.5), "Pourquoi soutenir cette version mobile", size=27, bold=True, color=COLORS["white"])
    add_text(
        slide,
        Inches(0.95),
        Inches(2.0),
        Inches(8.2),
        Inches(2.2),
        "Elle rapproche l'usager du soin, structure les alertes, facilite la supervision et transforme le smartphone en outil sanitaire de terrain pour les citoyens comme pour les institutions.",
        size=18,
        color=COLORS["white"],
    )
    add_card(slide, Inches(0.95), Inches(4.8), Inches(3.25), Inches(0.52), fill=COLORS["teal"], line=COLORS["teal"])
    add_text(slide, Inches(1.2), Inches(4.88), Inches(2.6), Inches(0.24), "Solution prete a demonstrer", size=13, bold=True, color=COLORS["white"])
    add_phone_frame(slide, SHOTS_DIR / "urgence-sanitaire.png", Inches(8.4), Inches(0.95), Inches(3.6), Inches(5.9))

    prs.save(OUTPUT)
    print(f"PPTX cree: {OUTPUT}")


if __name__ == "__main__":
    build()
